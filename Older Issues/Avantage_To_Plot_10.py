# -*- coding: utf-8 -*-
"""
Spyder Editor

This is a temporary script file.
"""

## Initialisation

#%pylab inline

from array import array
from pptx import Presentation
from pptx.util import Cm
from pptx.util import Pt

import numpy as np
import matplotlib.pyplot as plt
import time


from os import listdir
from os.path import isfile, join
import os, errno

# Import `os` 
import os

# Import pandas
import pandas as pd

#import tkinter as Tk
#from PyQt5 import QtCore, QtGui, QtWidgets


# List all files and directories in current directory
os.listdir('.')



#MAIN CONFIG******************************************************************
Config_filename = ''
for file in os.listdir('.'):    
    if file.endswith(".xlsx"):
        if file.startswith("ConfigXPS_10"):
            Config_filename = file
            print("Config File: ", Config_filename)
Config_file = pd.ExcelFile(Config_filename)


Main_Settings = Config_file.parse('Main Settings')
Div_Area_Val = [1,1,1,1,1,1,1,1,1,1,1,1,1,1]
i_TabSetting = 1

for i in range(0,len(Main_Settings)):
    if Main_Settings.ix[i,0] == "Config Tab Number":
        i_TabSetting = Main_Settings.ix[i,1]
        print('Config Tab to use: ', i_TabSetting)
    elif Main_Settings.ix[i,0] == "Folder":
        AnaFolder = Main_Settings.ix[i,1+i_TabSetting]
        #if AnaFolder.endswith("\\"): a = 1
        #else: AnaFolder = AnaFolder + '\\'
        AnaFolder2 = Main_Settings.ix[i+1,1+i_TabSetting]
        if AnaFolder2.endswith("\\"): a = 1
        else: AnaFolder2 = AnaFolder2 + '\\'
        AnaFolder = AnaFolder + AnaFolder2
        print('Folder name: ', AnaFolder)
    elif Main_Settings.ix[i,0] == 'Files' :
        file =  Main_Settings.ix[i,1+i_TabSetting].replace("\n","").replace("\t","").replace(" ","").split(",")
        print('File: ', file)
    elif Main_Settings.ix[i,0] == 'Peak line & value' :        
        plotline = Main_Settings.ix[i,1+i_TabSetting]
        plotlineS = plotline
        plotlineVal = plotline
    elif Main_Settings.ix[i,0] == 'Visible Y values' :
        YaxisVis = Main_Settings.ix[i,1+i_TabSetting]
    elif Main_Settings.ix[i,0] == 'Automatic Shirley' :
        AutoShirley = False
        AutoShirley = Main_Settings.ix[i,1+i_TabSetting]
    elif Main_Settings.ix[i,0] == 'Do quantification' :
        isQuant1 = Main_Settings.ix[i,1+i_TabSetting]
    elif Main_Settings.ix[i,0] == 'Do C-C correction' :
        IsCarbonComp = Main_Settings.ix[i,1+i_TabSetting]
    elif Main_Settings.ix[i,0] == 'Auto height comp' :
        IsHeightAuto = Main_Settings.ix[i,1+i_TabSetting]
    elif Main_Settings.ix[i,0] == 'Manual Height comp' :
        HeightComp = Main_Settings.ix[i,1+i_TabSetting]
    elif Main_Settings.ix[i,0] == 'Show Figure in Python' :
        UseSourcePython = Main_Settings.ix[i,1+i_TabSetting]
    elif Main_Settings.ix[i,0] == 'Do normalisation by val' :
        Div_Area_Check = Main_Settings.ix[i,1+i_TabSetting]
    elif Main_Settings.ix[i,0] == 'Normalisation value' :
        tmpDiv_Area_Val = Main_Settings.ix[i,1+i_TabSetting].replace("\n","").replace("\t","").replace(" ","").split(",")  
        for i2 in range(0,len(file)) : 
            if tmpDiv_Area_Val[i2] == "file":
                xl = pd.ExcelFile(AnaFolder+file[i2])
                File_Title = xl.parse('Titles')
                if(File_Title.ix[2,0] == 'Normalisation'): Div_Area_Val[i2] = File_Title.ix[2,1]
                else : Div_Area_Val[i2] = 1
            else : Div_Area_Val[i2] = int(tmpDiv_Area_Val[i2])
    elif Main_Settings.ix[i,0] == 'Save Images' :
        IsSaveFig = Main_Settings.ix[i,1+i_TabSetting]
    elif Main_Settings.ix[i,0] == 'DPI level' :
        dpiVal = Main_Settings.ix[i,1+i_TabSetting]
    elif Main_Settings.ix[i,0] == "Do a presentation":
        IsPresentation = Main_Settings.ix[i,1+i_TabSetting]    
    elif Main_Settings.ix[i,0] == "Filename":
        P_Fname = Main_Settings.ix[i,1+i_TabSetting]    
    elif Main_Settings.ix[i,0] == "Title":
        P_Title = Main_Settings.ix[i,1+i_TabSetting]    
    elif Main_Settings.ix[i,0] == "Template Filename":
        P_Template = Main_Settings.ix[i,1+i_TabSetting]
    elif Main_Settings.ix[i,0] == "Change Legend default location?":
        isChangeLocation = Main_Settings.ix[i,1+i_TabSetting]
    elif Main_Settings.ix[i,0] == "Legend position":
        Location = Main_Settings.ix[i,1+i_TabSetting]
    elif Main_Settings.ix[i,0] == "Legend Offset X":
        CompOffsetX = float(Main_Settings.ix[i,1+i_TabSetting])
    elif Main_Settings.ix[i,0] == "Legend Offset Y":
        CompOffsetY = float(Main_Settings.ix[i,1+i_TabSetting])        
    elif Main_Settings.ix[i,0] == "Autoheight for labelling":
        LabellingAutoHeight = Main_Settings.ix[i,1+i_TabSetting]
    elif Main_Settings.ix[i,0] == "File number for labelling":
        LabellingNumber = int(Main_Settings.ix[i,1+i_TabSetting]) - 1
        if LabellingNumber > len(file)-1: LabellingNumber = 0
    elif Main_Settings.ix[i,0] == "Plot settings sample number":
        PlotSettingNumber = int(Main_Settings.ix[i,1+i_TabSetting]) - 1  
        print('Setting Number: ', PlotSettingNumber)
    elif Main_Settings.ix[i,0] == "Lines settings sample number":
        LineSettingNumber = int(Main_Settings.ix[i,1+i_TabSetting]) - 1 
        print('Setting Number: ', LineSettingNumber)
    elif Main_Settings.ix[i,0] == "Offset plots":
        UseOffsetComp = Main_Settings.ix[i,1+i_TabSetting]    
    elif Main_Settings.ix[i,0] == "Offset Gap in Percentage":
        OffsetPercent = Main_Settings.ix[i,1+i_TabSetting]    
    elif Main_Settings.ix[i,0] == "Offset Maximum Percentage":
        OffsetPercent2 = Main_Settings.ix[i,1+i_TabSetting] 
    elif Main_Settings.ix[i,0] == "Offset in Survey":
        OffsetSurvey2 = Main_Settings.ix[i,1+i_TabSetting] 
#MAIN CONFIG******************************************************************

#LINE SETTINGS**********************************************
DEF_lines = []
Line_Settings = Config_file.parse('Line Settings')

for i in range(2,len(Line_Settings)):    
    if np.isnan(float(Line_Settings.ix[i,0+LineSettingNumber*4])) == True: 
        print('Not a number: ', Line_Settings.ix[i,1+LineSettingNumber*4])
        break
    #print(i, Line_Settings.ix[i,1+LineSettingNumber*4])
    if Line_Settings.ix[i,3+LineSettingNumber*4]:        
        DEF_lines.extend([Line_Settings.ix[i,0+LineSettingNumber*4], Line_Settings.ix[i,1+LineSettingNumber*4], 
                          Line_Settings.ix[i,2+LineSettingNumber*4]])
#LINE SETTINGS**********************************************

#print('Length',len(Line_Settings),len(DEF_lines))
#print('DEF LINES   ',DEF_lines)

#PLOT SETTINGS**********************************************
Plot_Settings = Config_file.parse('Plot Settings')

for i in range(0,len(Plot_Settings)):
    if Plot_Settings.ix[i,0+PlotSettingNumber*4] == 'Au4f' :
        MinMaxAu4f = Plot_Settings.ix[i,3+PlotSettingNumber*4]
        MinAu4f = Plot_Settings.ix[i,1+PlotSettingNumber*4]
        MaxAu4f = Plot_Settings.ix[i,2+PlotSettingNumber*4]

    if Plot_Settings.ix[i,0+PlotSettingNumber*4] == 'Pb4f' :
        MinMaxPb4f = Plot_Settings.ix[i,3+PlotSettingNumber*4]
        MinPb4f = Plot_Settings.ix[i,1+PlotSettingNumber*4]
        MaxPb4f = Plot_Settings.ix[i,2+PlotSettingNumber*4]
        
    elif Plot_Settings.ix[i,0+PlotSettingNumber*4] == 'Ce4d' :
        MinMaxCe4d = Plot_Settings.ix[i,3+PlotSettingNumber*4]
        MinCe4d = Plot_Settings.ix[i,1+PlotSettingNumber*4]
        MaxCe4d = Plot_Settings.ix[i,2+PlotSettingNumber*4]
    
    elif Plot_Settings.ix[i,0+PlotSettingNumber*4] == 'Sr3d' :
        MinMaxSr3d = Plot_Settings.ix[i,3+PlotSettingNumber*4]
        MinSr3d = Plot_Settings.ix[i,1+PlotSettingNumber*4]
        MaxSr3d = Plot_Settings.ix[i,2+PlotSettingNumber*4]
    
    elif Plot_Settings.ix[i,0+PlotSettingNumber*4] == 'Ce5d' :
        MinMaxCe5d = Plot_Settings.ix[i,3+PlotSettingNumber*4]
        MinCe5d = Plot_Settings.ix[i,1+PlotSettingNumber*4]
        MaxCe5d = Plot_Settings.ix[i,2+PlotSettingNumber*4]
    
    elif Plot_Settings.ix[i,0+PlotSettingNumber*4] == 'Ce3d' :
        MinMaxCe3d = Plot_Settings.ix[i,3+PlotSettingNumber*4]
        MinCe3d = Plot_Settings.ix[i,1+PlotSettingNumber*4]
        MaxCe3d = Plot_Settings.ix[i,2+PlotSettingNumber*4]
    
    elif Plot_Settings.ix[i,0+PlotSettingNumber*4] == 'La5d' :
        MinMaxLa5d = Plot_Settings.ix[i,3+PlotSettingNumber*4]
        MinLa5d = Plot_Settings.ix[i,1+PlotSettingNumber*4]
        MaxLa5d = Plot_Settings.ix[i,2+PlotSettingNumber*4]
 
    elif Plot_Settings.ix[i,0+PlotSettingNumber*4] == 'La4d' :
        MinMaxLa4d = Plot_Settings.ix[i,3+PlotSettingNumber*4]
        MinLa4d = Plot_Settings.ix[i,1+PlotSettingNumber*4]
        MaxLa4d = Plot_Settings.ix[i,2+PlotSettingNumber*4]  
    
    elif Plot_Settings.ix[i,0+PlotSettingNumber*4] == 'La3d' :
        MinMaxLa3d = Plot_Settings.ix[i,3+PlotSettingNumber*4]
        MinLa3d = Plot_Settings.ix[i,1+PlotSettingNumber*4]
        MaxLa3d = Plot_Settings.ix[i,2+PlotSettingNumber*4]
    
    elif Plot_Settings.ix[i,0+PlotSettingNumber*4] == 'Ag3d' :
        MinMaxAg3d = Plot_Settings.ix[i,3+PlotSettingNumber*4]
        MinAg3d = Plot_Settings.ix[i,1+PlotSettingNumber*4]
        MaxAg3d = Plot_Settings.ix[i,2+PlotSettingNumber*4]
    
    elif Plot_Settings.ix[i,0+PlotSettingNumber*4] == 'Zr3d' :
        MinMaxZr3d = Plot_Settings.ix[i,3+PlotSettingNumber*4]
        MinZr3d = Plot_Settings.ix[i,1+PlotSettingNumber*4]
        MaxZr3d = Plot_Settings.ix[i,2+PlotSettingNumber*4]

    elif Plot_Settings.ix[i,0+PlotSettingNumber*4] == 'Ru3d' :
        MinMaxRu3d = Plot_Settings.ix[i,3+PlotSettingNumber*4]
        MinRu3d = Plot_Settings.ix[i,1+PlotSettingNumber*4]
        MaxRu3d = Plot_Settings.ix[i,2+PlotSettingNumber*4]
    
    elif Plot_Settings.ix[i,0+PlotSettingNumber*4] == 'Sr3p' :
        MinMaxSr3p = Plot_Settings.ix[i,3+PlotSettingNumber*4]
        MinSr3p = Plot_Settings.ix[i,1+PlotSettingNumber*4]
        MaxSr3p = Plot_Settings.ix[i,2+PlotSettingNumber*4]

    elif Plot_Settings.ix[i,0+PlotSettingNumber*4] == 'Ru3p' :
        MinMaxRu3p = Plot_Settings.ix[i,3+PlotSettingNumber*4]
        MinRu3p = Plot_Settings.ix[i,1+PlotSettingNumber*4]
        MaxRu3p = Plot_Settings.ix[i,2+PlotSettingNumber*4]

    elif Plot_Settings.ix[i,0+PlotSettingNumber*4] == 'Mn2p' :
        MinMaxMn2p = Plot_Settings.ix[i,3+PlotSettingNumber*4]
        MinMn2p = Plot_Settings.ix[i,1+PlotSettingNumber*4]
        MaxMn2p = Plot_Settings.ix[i,2+PlotSettingNumber*4]
    
    elif Plot_Settings.ix[i,0+PlotSettingNumber*4] == 'Fe2p' :
        MinMaxFe2p = Plot_Settings.ix[i,3+PlotSettingNumber*4]
        MinFe2p = Plot_Settings.ix[i,1+PlotSettingNumber*4]
        MaxFe2p = Plot_Settings.ix[i,2+PlotSettingNumber*4]
    
    elif Plot_Settings.ix[i,0+PlotSettingNumber*4] == 'Co2p' :
        MinMaxCo2p = Plot_Settings.ix[i,3+PlotSettingNumber*4]
        MinCo2p = Plot_Settings.ix[i,1+PlotSettingNumber*4]
        MaxCo2p = Plot_Settings.ix[i,2+PlotSettingNumber*4]

    elif Plot_Settings.ix[i,0+PlotSettingNumber*4] == 'Cr2p' :
        MinMaxCr2p = Plot_Settings.ix[i,3+PlotSettingNumber*4]
        MinCr2p = Plot_Settings.ix[i,1+PlotSettingNumber*4]
        MaxCr2p = Plot_Settings.ix[i,2+PlotSettingNumber*4]
    
    elif Plot_Settings.ix[i,0+PlotSettingNumber*4] == 'Ni2p' :
        MinMaxNi2p = Plot_Settings.ix[i,3+PlotSettingNumber*4]
        MinNi2p = Plot_Settings.ix[i,1+PlotSettingNumber*4]
        MaxNi2p = Plot_Settings.ix[i,2+PlotSettingNumber*4]

    elif Plot_Settings.ix[i,0+PlotSettingNumber*4] == 'Ga2p' :
        MinMaxGa2p = Plot_Settings.ix[i,3+PlotSettingNumber*4]
        MinGa2p = Plot_Settings.ix[i,1+PlotSettingNumber*4]
        MaxGa2p = Plot_Settings.ix[i,2+PlotSettingNumber*4]

    elif Plot_Settings.ix[i,0+PlotSettingNumber*4] == 'Ga5p' :
        MinMaxGa5p = Plot_Settings.ix[i,3+PlotSettingNumber*4]
        MinGa5p = Plot_Settings.ix[i,1+PlotSettingNumber*4]
        MaxGa5p = Plot_Settings.ix[i,2+PlotSettingNumber*4]

    elif Plot_Settings.ix[i,0+PlotSettingNumber*4] == 'P2p' :
        MinMaxP2p = Plot_Settings.ix[i,3+PlotSettingNumber*4]
        MinP2p = Plot_Settings.ix[i,1+PlotSettingNumber*4]
        MaxP2p = Plot_Settings.ix[i,2+PlotSettingNumber*4]

    elif Plot_Settings.ix[i,0+PlotSettingNumber*4] == 'Si2p' :
        MinMaxSi2p = Plot_Settings.ix[i,3+PlotSettingNumber*4]
        MinSi2p = Plot_Settings.ix[i,1+PlotSettingNumber*4]
        MaxSi2p = Plot_Settings.ix[i,2+PlotSettingNumber*4]

    elif Plot_Settings.ix[i,0+PlotSettingNumber*4] == 'S2p' :
        MinMaxS2p = Plot_Settings.ix[i,3+PlotSettingNumber*4]
        MinS2p = Plot_Settings.ix[i,1+PlotSettingNumber*4]
        MaxS2p = Plot_Settings.ix[i,2+PlotSettingNumber*4]    

    elif Plot_Settings.ix[i,0+PlotSettingNumber*4] == 'Al2p' :
        MinMaxAl2p = Plot_Settings.ix[i,3+PlotSettingNumber*4]
        MinAl2p = Plot_Settings.ix[i,1+PlotSettingNumber*4]
        MaxAl2p = Plot_Settings.ix[i,2+PlotSettingNumber*4]  
    
    elif Plot_Settings.ix[i,0+PlotSettingNumber*4] == 'O1s' :
        MinMaxO1s = Plot_Settings.ix[i,3+PlotSettingNumber*4]
        MinO1s = Plot_Settings.ix[i,1+PlotSettingNumber*4]
        MaxO1s = Plot_Settings.ix[i,2+PlotSettingNumber*4]

    elif Plot_Settings.ix[i,0+PlotSettingNumber*4] == 'Na1s' :
        MinMaxNa1s = Plot_Settings.ix[i,3+PlotSettingNumber*4]
        MinNa1s = Plot_Settings.ix[i,1+PlotSettingNumber*4]
        MaxNa1s = Plot_Settings.ix[i,2+PlotSettingNumber*4]
    
    elif Plot_Settings.ix[i,0+PlotSettingNumber*4] == 'N1s' :
        MinMaxN1s = Plot_Settings.ix[i,3+PlotSettingNumber*4]
        MinN1s = Plot_Settings.ix[i,1+PlotSettingNumber*4]
        MaxN1s = Plot_Settings.ix[i,2+PlotSettingNumber*4]
    
    elif Plot_Settings.ix[i,0+PlotSettingNumber*4] == 'C1s' :
        MinMaxC1s = Plot_Settings.ix[i,3+PlotSettingNumber*4]
        MinC1s = Plot_Settings.ix[i,1+PlotSettingNumber*4]
        MaxC1s = Plot_Settings.ix[i,2+PlotSettingNumber*4]

    elif Plot_Settings.ix[i,0+PlotSettingNumber*4] == 'F1s' :
        MinMaxF1s = Plot_Settings.ix[i,3+PlotSettingNumber*4]
        MinF1s = Plot_Settings.ix[i,1+PlotSettingNumber*4]
        MaxF1s = Plot_Settings.ix[i,2+PlotSettingNumber*4]
    
    elif Plot_Settings.ix[i,0+PlotSettingNumber*4] == 'Valence' :
        MinMaxval = Plot_Settings.ix[i,3+PlotSettingNumber*4]
        Minval = Plot_Settings.ix[i,1+PlotSettingNumber*4]
        Maxval = Plot_Settings.ix[i,2+PlotSettingNumber*4]
#PLOT SETTINGS**********************************************

CarbonComp = [284.8, 284.8,284.8,284.8,284.8,284.8,284.8,284.8,284.8,284.8,284.8,284.8,284.8,284.8,284.8,284.8]


#print('is Area ', Div_Area_Check)

try:
    os.makedirs(AnaFolder+'PNGs')
except OSError as e:
    if e.errno != errno.EEXIST:
        raise

try:
    os.makedirs(AnaFolder+'Data_CSV')
except OSError as e:
    if e.errno != errno.EEXIST:
        raise


plotname =' '

def iter_cells(table):
    for row in table.rows:
        for cell in row.cells:
            yield cell

#-----FUNCTION FOR PLOTTING LINE-------------------------          
def plot_line(DEF_lines, x_BE,plotMinMax, plotMin,plotMax, TMP_y, y_An):
    TMP_lines = DEF_lines
    TMP_x_BE = x_BE
    x_line = []
    y_line = []
    lab_line = []
    lab_line_x = []
    change_line_x = []
    
    #print('Length tab', len(x_BE), '  ', len(TMP_y))
    

        
        
    if plotMinMax: 
        a=1
        for ix_Max in range (i_start,len(TMP_y)+i_start):
            if x_BE.ix[ix_Max] < plotMax: break
        for ix_Min in range (i_start,len(TMP_y)+i_start):
            if x_BE.ix[ix_Min] < plotMin: break
        I_min = min(TMP_y.ix[ix_Max:ix_Min])
        I_max = max(TMP_y.ix[ix_Max:ix_Min])
        #print('I_min: ',I_min, 'I_max: ', I_max)
    else:
        plotMin = min(TMP_x_BE)
        plotMax = max(TMP_x_BE)
        I_max = max(TMP_y)
        I_min = min(TMP_y)
    

    for m in range(0,int(len(TMP_lines)/3)):
        x_line.extend([TMP_lines[3*m],TMP_lines[3*m],TMP_lines[3*m]])
        y_line.extend([-500000, (y_An-0.04)*I_max,-500000])
        lab_line.extend([TMP_lines[3*m+1]])
        lab_line_x.extend([TMP_lines[3*m]])
        change_line_x.extend([TMP_lines[3*m+2]])
    
    #print(lab_line)
    for n in range(0,len(lab_line)):

         
        if lab_line_x[n] < plotMax:
            if lab_line_x[n] > plotMin:               
                for x_i in range(i_start,len(TMP_x_BE)):
                    if TMP_x_BE.ix[x_i]< lab_line_x[n]:
                        #print('found it    ',TMP_x_BE.ix[x_i])
                        break
                TMP_y_2 = TMP_y.ix[x_i-7:x_i+7]
                
                
                #New Line x_value
                if change_line_x[n] == True : 
                    New_X = round(x_BE.ix[np.argmax(TMP_y_2)],2)
                    if New_X == round(x_BE.ix[x_i],2) : print(lab_line[n], 'Same Value')
                    else : print(lab_line[n], ': ', 'Old: ', round(x_BE.ix[x_i],2), 'eV -->New: ', New_X, 'eV')
                    lab_line_x[n] = New_X
                    x_line[3*n] = New_X
                    x_line[3*n+1] = New_X
                    x_line[3*n+2] = New_X
                else : print(lab_line[n], 'Same Value: ', lab_line_x[n],'eV')
                plt.annotate(lab_line[n], xy=(lab_line_x[n], y_An*I_max), color = 'black', 
                     ha="center", va="bottom", rotation = 90, alpha = 0.8)
                if plotlineVal == True:
                    plt.annotate(lab_line_x[n], xy=(lab_line_x[n], I_min-0.01*I_max), color = 'lightgray', 
                     ha="right", va="bottom", rotation = 90)
    
    
    plt.plot(x_line,y_line,color = 'lightgray', linestyle='--', alpha = 0.4)

    return;
#--------------------------------------------------------
          
          
#-----FUNCTION FOR PLOTTING------------------------------
def plot_multi( tmp_fig, tmp_x, tmp_y, tmp_i,plotname, plotMinMax, x_Min, x_Max ):
    
    ax = tmp_fig.add_subplot(4,4,tmp_i+1)
    ax.yaxis.set_visible(YaxisVis)  
    #---------------------------------
    
    if plotMinMax == True:
        for ix_Max in range (i_start,len(tmp_y)+i_start):
            if tmp_x.ix[ix_Max] < x_Max: break
        for ix_Min in range (i_start,len(tmp_y)+i_start):
            if tmp_x.ix[ix_Min] < x_Min: break
        tmp_y = 100 * tmp_y / max(tmp_y.ix[ix_Max:ix_Min])
        I_min = min(tmp_y.ix[ix_Max:ix_Min])
        I_max = max(tmp_y.ix[ix_Max:ix_Min])
        if plotname == 'La 5d':
            x_Min = tmp_x.min()
            x_Max = tmp_x.max()
            I_min = min(tmp_y)
            I_max = tmp_y.max()                        
    else:
        tmp_y = 100 * tmp_y / max(tmp_y)
        I_min = min(tmp_y)
        I_max = tmp_y.max()
        x_Min = tmp_x.min()
        x_Max = tmp_x.max()
       
    ax.axis([x_Max, x_Min, I_min-I_max*0.05, I_max*1.55])


    #------PLOT CONFIGURATION---------
    plt.ylabel("")
    plt.xlabel("Binding Energy (eV)")
    plt.tight_layout()
    #---------------------------------     
               
    if plotname == 'Ce 5d': plotname = 'Ce 3d'     
    elif plotname == 'La 5d': plotname = 'Ce & La 3d'      
    
    
    #-----SET TITLE-------------------
    ax.text(0.01,.92,plotname, horizontalalignment = "left", transform=ax.transAxes, 
            fontsize=16, fontweight='bold'  )
    #---------------------------------
    plt.plot(tmp_x,tmp_y)#,'k')
    return(tmp_y);
#--------------------------------------------------------

#-----FUNCTION FOR PLOTTING------------------------------
def plot_settings( ):
    
    #ax = tmp_fig.add_subplot(4,4,tmp_i+1)
    ax.yaxis.set_visible(YaxisVis)  
    #---------------------------------
    
    #ax.axis([tmp_x.max()-0, tmp_x.min(), tmp_y.min()-tmp_y.max()*0.05, tmp_y.max()*1.2])


    #------PLOT CONFIGURATION---------
    plt.ylabel("")
    plt.xlabel("Binding Energy (eV)")
    plt.tight_layout()
    #---------------------------------            
                
    #-----SET TITLE-------------------
    ax.text(0.01,.92,plotname, horizontalalignment = "left", transform=ax.transAxes, 
            fontsize=16, fontweight='bold'  )
    #---------------------------------

    return;
#--------------------------------------------------------


#-----FUNCTION FOR PLOTTING COMP-------------------------
def plot_comp(TmpDef1, SmpName,TmpNumber,TmpName,UseMinMax,TmpMin,TmpMax,Tmpi):
    if len(TmpDef1.columns) > 2:  
        IsComp = True
        
        plt.figure(TmpNumber)
        plotname = TmpName
        if(TmpNumber ==4): ax = fig_multi.add_subplot(4,4,i2+1)
        if(TmpNumber ==5): ax = fig_multiS.add_subplot(4,4,i2s+1)
        if(TmpNumber ==7): ax = plt.gca()
        if(TmpNumber ==11): ax = plt.gca()
        if(TmpNumber ==10):             
            plt.figure(10)
            ax = plt.gca()
            #fig_survey, ax = plt.gca()
            #plt.figure(3, figsize=(10,5))
            #ax = plt.gca()
            #plt.figure(10, figsize=(10,5))    
            offset1 = -0.05 #-0.15
        else: offset1 = 0
        ax.yaxis.set_visible(YaxisVis) 
        plt.ylabel("")
        plt.xlabel("Binding Energy (eV)")
        plt.tight_layout()
        ax.text(0.01,.92,plotname, horizontalalignment = "left", transform=ax.transAxes, 
            fontsize=16, fontweight='bold'  )
        I_max =0
        I_min = 100
        
        '''
        I_maxCOMP = 0        
        for i in range(0, int(len(TmpDef1.columns)/2)):
            tmplist = list(TmpDef1)
            for ix_Max in range (i_start,len(TmpDef1[tmplist[2*i]])+i_start):
                if TmpDef1[tmplist[2*i]].ix[ix_Max] < TmpMax: break
            for ix_Min in range (i_start,len(TmpDef1[tmplist[2*i]])+i_start):
                if TmpDef1[tmplist[2*i]].ix[ix_Min] < TmpMin: break

            if I_maxCOMP <  max(TmpDef1[tmplist[2*i+1]].ix[ix_Max:ix_Min]):
                I_maxCOMP = max(TmpDef1[tmplist[2*i+1]].ix[ix_Max:ix_Min])
                print("I_max: ", I_maxCOMP)
        '''
            
        for i in range(0, int(len(TmpDef1.columns)/2)):
            tmplist = list(TmpDef1)
            if plotname == 'Ce & La 3d':
                X_Min = TmpDef1[tmplist[2*i]].min()+1
                X_Max = TmpDef1[tmplist[2*i]].max()-1
                for ix_Max in range (i_start,len(TmpDef1[tmplist[2*i]])+i_start):
                    if TmpDef1[tmplist[2*i]].ix[ix_Max] < TmpMax: 
                        break
                for ix_Min in range (i_start,len(TmpDef1[tmplist[2*i]])+i_start):
                    if TmpDef1[tmplist[2*i]].ix[ix_Min] < TmpMin: break
                if I_min > min(TmpDef1[tmplist[2*i+1]].ix[ix_Max:ix_Min]):
                    I_min = min(TmpDef1[tmplist[2*i+1]].ix[ix_Max:ix_Min])                
                #I_max = max(TmpDef1[tmplist[2*i+1]].ix[ix_Max:ix_Min])  
                if I_max < max(TmpDef1[tmplist[2*i+1]]) : I_max = max(TmpDef1[tmplist[2*i+1]])
            elif UseMinMax == True:
                X_Min = TmpMin
                X_Max = TmpMax
                for ix_Max in range (i_start,len(TmpDef1[tmplist[2*i]])+i_start):
                    if TmpDef1[tmplist[2*i]].ix[ix_Max] < TmpMax: break
                for ix_Min in range (i_start,len(TmpDef1[tmplist[2*i]])+i_start):
                    if TmpDef1[tmplist[2*i]].ix[ix_Min] < TmpMin: break
                if I_min > min(TmpDef1[tmplist[2*i+1]].ix[ix_Max:ix_Min]):
                    I_min = min(TmpDef1[tmplist[2*i+1]].ix[ix_Max:ix_Min])                
                if I_max <  max(TmpDef1[tmplist[2*i+1]].ix[ix_Max:ix_Min]):
                    I_max = max(TmpDef1[tmplist[2*i+1]].ix[ix_Max:ix_Min])
                
            else:
                X_Min = TmpDef1[tmplist[2*i]].min()+1
                X_Max = TmpDef1[tmplist[2*i]].max()-1
                if TmpDef1[tmplist[2*i+1]].max() > I_max: 
                    I_max = TmpDef1[tmplist[2*i+1]].max()
                if TmpDef1[tmplist[2*i+1]].min() < I_min: 
                    I_min = TmpDef1[tmplist[2*i+1]].min()

            #######OFFSET#########################################
            offsetSurvey = 0.4*(I_max-I_min)
            if UseOffsetComp == True and plotname != 'Survey':
                if i == 0: 
                    offset1 = (I_max - I_min)*OffsetPercent
                    offset2 = (I_max - I_min)*OffsetPercent2
                    #print("offset Made")
                else: a = 0
                
            elif(TmpNumber == 10): OffsetSurvey = OffsetSurvey2*(I_max-I_min)
            #######OFFSET#########################################           
            
            #plt.plot(TmpDef1[tmplist[2*i]],TmpDef1[tmplist[2*i+1]]+offset1*i*100, label= SmpName[i])
            if(TmpNumber == 10): plt.plot(TmpDef1[tmplist[2*i]],TmpDef1[tmplist[2*i+1]]+OffsetSurvey*i, label= SmpName[i])
            else: plt.plot(TmpDef1[tmplist[2*i]],TmpDef1[tmplist[2*i+1]]+offset1*i, label= SmpName[i])
            #plt.legend(loc=1,frameon=False)
            #print(Location)
            
            if isChangeLocation == True : plt.legend(loc=Location, bbox_to_anchor=(CompOffsetX, CompOffsetY),frameon=False)
            else : plt.legend(loc='upper right', bbox_to_anchor=(1, 1),frameon=False)
            #plt.draw_frame(False)
            if(TmpNumber ==10):                         
                plt.axis([X_Max, X_Min, I_min-(I_max-I_min)*0.05+OffsetSurvey*int(len(TmpDef1.columns)/2-1), I_max*(1)*1.25])
                #plt.axis([X_Max, X_Min, I_min-(I_max-I_min)*0.05 + 100*offset1*i, I_max*(1)*1.25])

                if i == 0:                    
                    i4 =10
                       
                    for k in range(1, len(survey_label)+1):                        
                        for l in range(i_start, len(survey_BE)):
                            i4 = l
                            if float(survey_BE[l]) < float(survey_position[k]):
                                i4 = l
                                break
                        TMP_survey_y = TmpDef1[tmplist[2*i+1]]
                        survey_yMAX = TMP_survey_y[i4-30:i4+10]
                        #survey_yMAX = survey_y[i4-30:i4+10]
                                                
                        ax.annotate(survey_label[k],xy=(survey_position[k], int(2+max(survey_yMAX))),ha="center", va="bottom", 
                                         rotation = 90)#, xytext=(survey_position[k], int(10+max(survey_yMAX))))
                
            elif i == 0 and UseOffsetComp == True:
                if offset1 < 0:
                    #print("offset Made2")
                    plt.axis([X_Max, X_Min, I_min-((I_max-I_min)*0.05)+offset1*int(len(TmpDef1.columns)/2-1), I_max* HeightComp+offset2]) 
            elif UseOffsetComp == False: plt.axis([X_Max, X_Min, I_min-((I_max-I_min)*0.05), I_max* HeightComp])  #(1+offset1)*HeightComp])
        Tmpi = Tmpi + 1  
    
    return(Tmpi);
#--------------------------------------------------------

#-----FUNCTION FOR PLOTTING FIT--------------------------
def plot_fit(x_fit, data_fit,TmpName, TmpNumber,plotMinMax, plotMin,plotMax,i_fit ):
    if len(data_fit.columns) > 6:  
        col = ['C1','C2','C3','C4','C5','C6','C7','C8','C9','C10']
        i_c = 0
        #x_fit = data_fit.ix[i_start:,0]
        y_fit = data_fit.ix[i_start:,j_start+2]
        if data_fit.ix[i_start-2,j_start+4] == 'Backgnd.': 
            y_BKG = data_fit.ix[i_start:,j_start+4]
            i_0 = j_start+5
        else:
            i_0 = j_start+4
            for i_shirley in range(0, len(df1.columns)):
                if df1.ix[i_start-2,i_shirley] == 'Backgnd.':
                    y_BKG = data_fit.ix[i_start:,i_shirley]
        
        X_Min = x_fit.min()
        X_Max = x_fit.max()
        I_min = y_fit.min()
        I_max = y_fit.max()
        
        if plotMinMax == True:
            for ix_Max in range (i_start,len(y_fit)+i_start):
                if x_fit.ix[ix_Max] < plotMax: break
            for ix_Min in range (i_start,len(y_fit)+i_start):
                if x_fit.ix[ix_Min] < plotMin: break
            #y_fit = 100 * y_fit / max(y_fit.ix[ix_Max:ix_Min])
            I_min = min(y_fit.ix[ix_Max:ix_Min])
            I_max = max(y_fit.ix[ix_Max:ix_Min])
        
        plt.figure(TmpNumber)
        plotname = TmpName
        if(TmpNumber ==8): ax = fig_fit.add_subplot(2,4,i_fit+1)
        ax.yaxis.set_visible(YaxisVis) 
        plt.ylabel("")
        plt.xlabel("Binding Energy (eV)")
        plt.tight_layout()
        ax.text(0.01,.92,plotname, horizontalalignment = "left", transform=ax.transAxes, 
            fontsize=16, fontweight='bold'  )
        x_f = x_fit.values
        y_b = y_BKG.values
        y_f = y_fit.values
        #print(y_b)
        plt.plot(x_f,y_f, marker = 'o', markerfacecolor= 'white' )
        plt.plot(x_f,y_b,color = 'lightgray')#, linestyle='--')
        for i in range(i_0, int(len(data_fit.columns))):            
            label_fit = data_fit.ix[i_start-2,i]
            if label_fit == 'Residuals': break
            y_fit = data_fit.ix[i_start:,i]
            y_f = y_fit.values
            if label_fit == 'Envelope':
                 plt.plot(x_f,y_f, 'k')
            elif label_fit =='Backgnd.':
                plt.plot(x_f,y_f, color = 'lightgray', linestyle='--')     
                a=1
                   
            #plt.fill_between(x_f,y_f,y_b)#, interpolate =True)#,facecolor='grey', alpha=0.5, interpolate=True)
            else:
                if label_fit == '-': i_c=i_c-1
                plt.plot(x_f,y_f,col[i_c])#'k',linestyle='--')# label=label_fit)
                #print(label_fit,'   ', col[i_c])                

                #print(x_f_line)
                if label_fit == '-':
                    a =1
                    
                else:
                    y_f2 = y_f-y_b
                    x_f_line = x_f[np.argmax(y_f2)]
                    if x_f_line > plotMin:
                        if x_f_line < plotMax:                            
                            plt.plot([x_f_line,x_f_line],[-I_max,I_max+(I_max-I_min)*0.1],col[i_c],  linestyle='--', alpha = 0.4) #color = 'lightgray', linestyle='--')
                            plt.annotate(label_fit, xy=(x_f_line, I_max+(I_max-I_min)*.15), color = col[i_c],
                                         ha="center", va="bottom", rotation = 90)  #color = 'grey',
                i_c=i_c+1
        
        plt.legend(loc=1)
        plt.axis([plotMax, plotMin, 
                      I_min-(I_max-I_min)*0.05, I_max+(I_max-I_min)*0.6])
    
    return();
#--------------------------------------------------------

#-----FUNCTION FOR PLOTTING ETCH-------------------------
def plot_etch(x_fit, data_fit,TmpName, TmpNumber,plotMinMax, plotMin,plotMax,i_fit ):
    if len(data_fit.columns) > 6:  
        x_fit = data_fit.ix[i_start:,j_start+0]
        y_fit = data_fit.ix[i_start:,j_start+2]
        X_Min = x_fit.min()
        X_Max = x_fit.max()
        I_min = y_fit.min()
        I_max = y_fit.max()
        
        if plotMinMax == True:
            for ix_Max in range (i_start,len(y_fit)+i_start):
                if x_fit.ix[ix_Max] < plotMax: break
            for ix_Min in range (i_start,len(y_fit)+i_start):
                if x_fit.ix[ix_Min] < plotMin: break
            I_min = min(y_fit.ix[ix_Max:ix_Min])
            I_max = max(y_fit.ix[ix_Max:ix_Min])
        SpaceEtch = 0.1*(I_max-I_min)
        plt.figure(TmpNumber)
        plotname = TmpName
        if(TmpNumber ==9): 
            ax = fig_etch.add_subplot(2,4,i_fit+1)
            ax.text(0.01,.95,plotname, horizontalalignment = "left", transform=ax.transAxes, 
                    fontsize=16, fontweight='bold'  )            
        if(TmpNumber ==10): 
            plt.figure(10)#, figsize=(10,7))
            ax = plt.gca()
            ax.text(0.01,.91,plotname, horizontalalignment = "left", transform=ax.transAxes, 
                    fontsize=16, fontweight='bold'  )             
        ax.yaxis.set_visible(YaxisVis) 
        plt.ylabel("")
        plt.xlabel("Binding Energy (eV)")
        plt.tight_layout()

        x_f = x_fit.values
        for i in range(2, int(len(data_fit.columns))):
            y_fit = data_fit.ix[16:,i]
            y_f = y_fit.values +SpaceEtch*(i-2)
            tmpMax = max(y_f)
            tmpMin = min(y_f)
            if tmpMax > I_max: I_max = tmpMax
            if tmpMin < I_min: I_min = tmpMin
            
            plt.plot(x_f,y_f)        
        
        plt.legend(loc=1)
        plt.axis([plotMax, plotMin, 
                      I_min-(I_max-I_min)*0.05, I_max+(I_max-I_min)*0.1])
    
    return();
#--------------------------------------------------------


df_survey = pd.DataFrame()
df_survey2 = pd.DataFrame()
df_val = pd.DataFrame()
df_Ce5d = pd.DataFrame()
df_Au4f = pd.DataFrame()
df_Pb4f = pd.DataFrame()
df_Gd4d = pd.DataFrame()
df_Ce4d = pd.DataFrame()
df_La4d = pd.DataFrame()
df_Ce3d = pd.DataFrame()
df_La5d = pd.DataFrame()
df_La3d = pd.DataFrame()
df_Ag3d = pd.DataFrame()
df_Sr3d = pd.DataFrame()
df_Zr3d = pd.DataFrame()
df_Ru3d = pd.DataFrame()
df_Ru3p = pd.DataFrame()
df_Sr3p = pd.DataFrame()
df_Fe2p = pd.DataFrame()
df_Mn2p = pd.DataFrame()
df_Cr2p = pd.DataFrame()
df_S2p = pd.DataFrame()
df_Al2p = pd.DataFrame()
df_Si2p = pd.DataFrame()
df_Co2p = pd.DataFrame()
df_Ni2p = pd.DataFrame()
df_Ga5p = pd.DataFrame()
df_Ga2p = pd.DataFrame()
df_P2p = pd.DataFrame()
df_O1s = pd.DataFrame()
df_C1s = pd.DataFrame()
df_F1s = pd.DataFrame()
df_Na1s = pd.DataFrame()
df_N1s = pd.DataFrame()

df_Ce5ds = pd.DataFrame()
df_Au4fs = pd.DataFrame()
df_Pb4fs = pd.DataFrame()
df_Ce4ds = pd.DataFrame()
df_La4ds = pd.DataFrame()
df_Ce3ds = pd.DataFrame()
df_La5ds = pd.DataFrame()
df_La3ds = pd.DataFrame()
df_Ag3ds = pd.DataFrame()
df_Sr3ds = pd.DataFrame()
df_Zr3ds = pd.DataFrame()
df_Ru3ds = pd.DataFrame()
df_Ru3ps = pd.DataFrame()
df_Sr3ps = pd.DataFrame()
df_Fe2ps = pd.DataFrame()
df_Mn2ps = pd.DataFrame()
df_Cr2ps = pd.DataFrame()
df_S2ps = pd.DataFrame()
df_Al2ps = pd.DataFrame()
df_Si2ps = pd.DataFrame()
df_Co2ps = pd.DataFrame()
df_Ni2ps = pd.DataFrame()
df_Ga5ps = pd.DataFrame()
df_Ga2ps = pd.DataFrame()
df_P2ps = pd.DataFrame()
df_O1ss = pd.DataFrame()
df_C1ss = pd.DataFrame()
df_F1ss = pd.DataFrame()
df_Na1ss = pd.DataFrame()
df_N1ss = pd.DataFrame()

df_surveyT = []
df_surveyT2 = []
df_valT = []
df_Ce5dT = []
df_Au4fT = []
df_Pb4fT = []
df_Gd4dT = []
df_Ce4dT = []
df_La4dT = []
df_Ce3dT = []
df_La5dT = []
df_La3dT = []
df_Ag3dT = []
df_Sr3dT = []
df_Zr3dT = []
df_Ru3dT = []
df_Ru3pT = []
df_Sr3pT = []
df_Fe2pT = []
df_Mn2pT = []
df_Cr2pT = []
df_S2pT = []
df_Al2pT = []
df_Si2pT = []
df_Co2pT = []
df_Ni2pT = []
df_Ga5pT = []
df_Ga2pT = []
df_P2pT = []
df_O1sT = []
df_C1sT = []
df_F1sT = []
df_Na1sT = []
df_N1sT = []

df_Ce5dsT = []
df_Au4fsT = []
df_Pb4fsT = []
df_Ce4dsT = []
df_La4dsT = []
df_Ce3dsT = []
df_La5dsT = []
df_La3dsT = []
df_Ag3dsT = []
df_Sr3dsT = []
df_Zr3dsT = []
df_Ru3dsT = []
df_Ru3psT = []
df_Sr3psT = []
df_Fe2psT = []
df_Mn2psT = []
df_Cr2psT = []
df_Al2psT = []
df_Si2psT = []
df_S2psT = []
df_Co2psT = []
df_Ni2psT = []
df_Ga5psT = []
df_Ga2psT = []
df_P2psT = []
df_O1ssT = []
df_C1ssT = []
df_F1ssT = []
df_Na1ssT = []
df_N1ssT = []

#INIT
plotname = ''
i1 = 0
i2 = 0
i2s = 0
i_fit =0
i_start = 14
j_start = 0

Sample_name=[]

isQuant = False
IsShirley = True
IsValence = False
IsComp = False
IsFit = False
IsEtch = False

#Presentation
if IsSaveFig == True:
    prs = Presentation(P_Template)
    prs.save(AnaFolder+P_Fname)
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]    
    
    title.text = P_Title
    subtitle.text = 'Gwilherm Kerherve, '+time.strftime("%d/%m/%Y")

    prs.save(AnaFolder+P_Fname)



#-----SINGLE PLOT----------------------------------------

for j in range(0, len(file)):
    IsShirley = False
    IsValence = False
    IsFit = False
    IsEtch = False    
    print('Filename: '+ file[j])
    
    # Load spreadsheet
    xl = pd.ExcelFile(AnaFolder+file[j])
    
    
    #Sample name 
    File_Title = xl.parse('Titles')
    Sample_name.append(File_Title.ix[3,0])
    print('Sample name: '+ File_Title.ix[3,0])
    
    if(File_Title.ix[4,0] == 'Carbon BE'): CarbonComp[j] = File_Title.ix[4,1]
    print('Carbon peak: ', CarbonComp[j], ' eV')
    #print(CarbonComp[j])
    
    #QUANTIFICATION TABLES----------------------------------------------------------------------------------------------------
    if isQuant1 == True:
        #print(xl.sheet_names)
        if ['Quantifications' in xl.sheet_names] : 
            print('Data quantified')
            isQuant = True
            File_Quant = xl.parse('Quantifications')
            #print(CarbonComp[j])
            #Tmp_col = File_Quant.columns.get_loc('Peak BE')

            Peak_Name = File_Quant.ix[:,File_Quant.columns.get_loc('Name ')]
            Peak_Name = Peak_Name[Peak_Name.notnull()]
            #Quant_Tab1 = File_Quant.ix[0:len(Peak_Name)-1,[File_Quant.columns.get_loc('Name '),
                                       #File_Quant.ix[0:len(Peak_Name)-1,File_Quant.columns.get_loc('Peak BE')], File_Quant.columns.get_loc('FWHM eV'), File_Quant.columns.get_loc('Atomic %')]]
            Quant_Tab1 = File_Quant.ix[0:len(Peak_Name)-1,[File_Quant.columns.get_loc('Name '),
                                       File_Quant.columns.get_loc('Peak BE'), File_Quant.columns.get_loc('FWHM eV'), File_Quant.columns.get_loc('Atomic %')]]
            Area_Conc = File_Quant.ix[0:len(Peak_Name)-1,File_Quant.columns.get_loc('Atomic %')]
            
            for j_tab in range(0, len(Peak_Name)):
                Quant_Tab1.ix[j_tab,1] = Quant_Tab1.ix[j_tab,1] - CarbonComp[j]+284.8
            print(Quant_Tab1)
            #print('row ',len(Peak_Name))
            #print('column  ', len(list(Quant_Tab1)))
            Column_Num1 = len(list(File_Quant))
            #print(Area_Conc)
            #print(Column_Num1)
                        
            if Column_Num1 > 6:
                Quant_Tab2 = File_Quant.ix[0:len(Peak_Name)-1, File_Quant.columns.get_loc('Atomic %')+1:]
                #print(Quant_Tab2)
                Header_Tab2 = list(Quant_Tab2)
                Quant_Tab2_Check = Quant_Tab2.notnull()
                
                #Quant_Tab3 =array([][])
                word1 = [' '] * (len(Header_Tab2)+1)
                word2 = [' '] * (len(Header_Tab2)+1)
                word1[0] = 'Name'
                word2[0] = 'Atomic Ratio'
                #print(word1)
                #print(word2)
                print(Header_Tab2)
                #Create concentration tables
                for i_tab in range(0 , len(Header_Tab2)):
                    list1 = list(Peak_Name[Quant_Tab2_Check.ix[:,i_tab]])
                    word1[i_tab+1] = list1[0]
                    for j_tab in range(1,len(list1)) : word1[i_tab+1] = word1[i_tab+1]+' : '+list1[j_tab]
                    #print(word1)
                    #Quant_Tab3.ix[i_tab+1,0] = word1
                    
                    
                    list2 = list(Area_Conc[Quant_Tab2_Check.ix[:,i_tab]])
                    list2 = np.round(list2 / sum(list2[:]) *100)
                    word2[i_tab+1] = str(list2[0])                    
                    for j_tab in range(1,len(list2)) : word2[i_tab+1] = word2[i_tab+1]+' : '+str(list2[j_tab])
                    #print(word2)
                    #Quant_Tab3[i_tab+1,1] = word2
                Quant_Tab3 = np.column_stack((word1,word2))
                
                print('Ratios obtained from Quantifications')
                for i_tab in range(0, len(Header_Tab2)): print(Quant_Tab3[i_tab][0] ,'       | ' , Quant_Tab3[i_tab][1])
                
        #END QUANTIFICATION TABLES----------------------------------------------------------------------------------------------
        else: 
            print('Data not quantified!!')
            isQuant = False
            isQuant1 = False
    
    # Set the sub figures
    fig = plt.figure(1, figsize=(14,14))
    fig_shirley1 = plt.figure(2, figsize=(14,14))
    fig_val = plt.figure(6, figsize=(4,4))
    fig_fit = plt.figure(8, figsize=(14,7))
    fig_etch = plt.figure(9, figsize=(14,14))
    
    plt.figure(1)
    # Reset indices
    i1 = 0    
    i3 = 0
    i5 = 0
    i_fit = 0
    i_etch = 0
    
    #Description of graphs
    print('Survey, Core levels and Valence band')
    


    # Load a sheet into a DataFrame by name: df1
    for i in range(0, len(xl.sheet_names)-1):    
        if(xl.sheet_names[i] == 'Peak Table'): #CHECK WHEN TO STOP LOOP
            i3 = i
            break
        
        df1 = xl.parse(xl.sheet_names[i])
        #print(df1)
        plotname = xl.sheet_names[i]
        plotname = plotname[0:2]+" "+plotname[2:4]
        if plotname[0:2] == 'O1':
            plotname = 'O 1s'
        if plotname[0:2] == 'C1':
            plotname = 'C 1s'        
        if plotname[0:2] == 'F1':
            plotname = 'F 1s' 
        if plotname[0:2] == 'N1':
            plotname = 'N 1s'
        if plotname[0:2] == 'S2':
            plotname = 'S 2p'   
        if plotname[0:2] == 'P2':
            plotname = 'P 2p'  
        if plotname[0:2] == 'Su':
            plotname = 'Survey'
        if plotname[0:2] == 'Va':
            plotname = 'Valence'        
        
        #print('GET LOC: ',df1.columns)#.get_loc('Binding Energy (E)'))
        

            
        j_start = 0
        for j_start in range(0, 4):
            #print(df1)
            x_BE = df1.ix[0:,j_start]
            isJ_Start = False
            for i_start in range(0, len(x_BE)):                
                if x_BE.ix[i_start] == 'eV':
                    isJ_Start = True
                    i_start = i_start+1
                    break        
            if isJ_Start: 
                break
        
        #print('I & J Start:  ',i_start, '  ', j_start)
        x_BE = df1.ix[i_start:,j_start]                
        
        TMPCarbonComp = 284.8 - CarbonComp[j]
        x_BE = x_BE+ TMPCarbonComp
        
        y = df1.ix[i_start:,j_start+2]
        
        plotMin = min(x_BE)
        plotMax = max(x_BE)
        plotMinMax = False
        
        if i_start-9 > 0:
            checkFit = df1.ix[i_start-9,j_start+1]
            if checkFit == 'Etch Time': 
                IsEtch = True
                #print(plotname)
                #print('...Etching...')        
        
        
        #-------------REMOVE BACKGROUND----------------------
        if plotname == 'Valence': y_BKG0 = sum(y[len(y)-5:len(y)-1])/4
        else : y_BKG0 = sum(y[len(y)-20:len(y)-10])/10
        #y_BKG0 = 10
        #print(plotname, '   ', y_BKG0)
        #print(y[len(y)-20:len(y)-10])
        #if(plotname != 'Survey'): y1 = y - y_BKG0
        #else: y1 = y
        y1 = y - y_BKG0
        if plotname == 'O 1s': print('O1s DIV AREA ==', Div_Area_Val[j])
        if plotname == 'Valence': y2 = y1/max(y1)
        elif Div_Area_Check == True: y2 = y1 / Div_Area_Val[j]
        else: 
            y2 = y1/max(y1)*100
            #y2 = y2/sum(y2)*100
        #----------------------------------------------------
        
        #-------------REMOVE SHIRLEY-------------------------
        if (AutoShirley == True) or (len(df1.columns) > j_start+4):
            if AutoShirley == True:
                IsShirley = True
                
                tmpI1 = y[i_start]
                tmpI2 = y[len(y)+i_start-1]    
                tmpA1A2 = sum(y)
                y_BKG1 = y+0.1
                for i_shirley in range(i_start, len(y)+i_start-1) :
                    tmpA1 = float(sum(y[i_shirley: len(y)+i_start-1]))
                    y_BKG1[i_shirley] = float(tmpI2 + (tmpI1 - tmpI2)*tmpA1/tmpA1A2)
                y3 = y - y_BKG1
                if Div_Area_Check == True: y3 = y3 / Div_Area_Val[j]
                else:
                    y3 = y3 / max(y3)
                    #y3 = y3 / sum(y3)*100
                
                y4 = 100*y3 / max(y3)                
            
            elif df1.ix[i_start-2,j_start+4] == 'Backgnd.':                
                IsShirley = True
                y_BKG1 = df1.ix[i_start:,j_start+4]
                y3 = y - y_BKG1
                if Div_Area_Check == True: y3 = y3 / Div_Area_Val[j]
                else:
                    y3 = y3 / max(y3)
                    #y3 = y3 / sum(y3)*100
                
                y4 = 100*y3 / max(y3)                       #MAY NEED TO TRANSFORM Y3 INTO Y4         
            elif len(df1.columns) > j_start+5:
                for i_shirley in range(0, len(df1.columns)):
                    if df1.ix[i_start-2,i_shirley] == 'Backgnd.':
                        print('GOT SHIRLEY')
                        IsShirley = True
                        y_BKG1 = df1.ix[i_start:,i_shirley]
                        y3 = y - y_BKG1
                        if Div_Area_Check == True: y3 = y3 / Div_Area_Val[j]
                        else:
                            y3 = y3 / max(y3)
                            #y3 = y3 / sum(y3)*100
                        
                        y4 = 100*y3 / max(y3)                       #MAY NEED TO TRANSFORM Y3 INTO Y4                            
        #----------------------------------------------------        
        else: y3 = y2     
        
        #-----WHICH CORE LEVEL FOR COMPARISON-------------------------
        #-----CHECK VALENCE----------------------------------
        if plotname == 'Valence':
            IsValence = True
            df_val['BE'+Sample_name[j]] = x_BE
            df_val[file[j]] = y2
            df_valT.append(Sample_name[j])

        #----------------------------------------------------

        #-----CHECK CE5D-------------------------------------
        if plotname == 'Ce 5d':
            if MinMaxCe5d == True:                                              #START HERE
                plotMinMax = True                           
                for i_Max in range (i_start,len(x_BE)+i_start):
                    if x_BE[i_Max] < MaxCe5d: break#                        HERE
                for i_Min in range (i_start,len(x_BE)+i_start):
                    if x_BE[i_Min] < MinCe5d: break#                        HERE       
                plotMax = x_BE[i_Max]
                plotMin = x_BE[i_Min]
                if Div_Area_Check == False:        
                    y_BKG = (y[i_Min-3]+y[i_Min-2]+y[i_Min-1]+y[i_Min])/4       
                    y2 = y - y_BKG                 
                    y2 = 100* y2 / max(y2[i_Max:i_Min-1])             
                    y3 = 100* y3 / max(y3[i_Max:i_Min-1])                       #FINISH HERE    
                else:
                    y_BKG = (y[i_Min-3]+y[i_Min-2]+y[i_Min-1]+y[i_Min])/4       
                    y2 = (y - y_BKG) / Div_Area_Val[j]       
            additional = pd.DataFrame({'BE'+Sample_name[j]: x_BE, file[j]:y2})
            new = pd.concat([df_Ce5d,additional],axis =1)
            df_Ce5d = new
            df_Ce5dT.append(Sample_name[j])
            if (AutoShirley == True) or (len(df1.columns) > j_start + 4):                
                if (AutoShirley == True) or (df1.ix[i_start-2,j_start+4] == 'Backgnd.'):
                    additional = pd.DataFrame({'BE'+Sample_name[j]: x_BE, file[j]:y3})
                    new = pd.concat([df_Ce5ds,additional],axis =1)
                    df_Ce5ds = new
                    df_Ce5dsT.append(Sample_name[j])
        #----------------------------------------------------

        #-----CHECK Au4F-------------------------------------
        if plotname == 'Au 4f':
            if MinMaxAu4f == True:                                              #START HERE
                plotMinMax = True                                   
                for i_Max in range (i_start,len(x_BE)+i_start):
                    if x_BE[i_Max] < MaxAu4f: break#                        HERE
                for i_Min in range (i_start,len(x_BE)+i_start):
                    if x_BE[i_Min] < MinAu4f: break#                        HERE       
                plotMax = x_BE[i_Max]
                plotMin = x_BE[i_Min]
                if Div_Area_Check == False: 
                    y_BKG = (y[i_Min-3]+y[i_Min-2]+y[i_Min-1]+y[i_Min])/4       
                    y2 = y - y_BKG                 
                    y2 = 100* y2 / max(y2[i_Max:i_Min-1])             
                    y3 = 100* y3 / max(y3[i_Max:i_Min-1])                       #FINISH HERE    
                else:
                    y_BKG = (y[i_Min-3]+y[i_Min-2]+y[i_Min-1]+y[i_Min])/4       
                    y2 = (y - y_BKG) / Div_Area_Val[j]       
            #df_Au4f['BE'+Sample_name[j]] = x_BE
            additional = pd.DataFrame({'BE'+Sample_name[j]: x_BE, file[j]:y2})
            new = pd.concat([df_Au4f,additional],axis =1)
            df_Au4f = new
            df_Au4fT.append(Sample_name[j])
            #print(df_Au4f.head())
            #df_Au4f[file[j]] = y2
            if (AutoShirley == True) or (len(df1.columns) > j_start + 4):                
                if (AutoShirley == True) or (df1.ix[i_start-2,j_start+4] == 'Backgnd.'):
                    #df_Au4fs['BE'+Sample_name[j]] = x_BE
                    #df_Au4fs[file[j]] = y3
                    additional = pd.DataFrame({'BE'+Sample_name[j]: x_BE, file[j]:y3})
                    new = pd.concat([df_Au4fs,additional],axis =1)
                    df_Au4fs = new
                    df_Au4fsT.append(Sample_name[j])
        #----------------------------------------------------

        #-----CHECK PB4F-------------------------------------
        if plotname == 'Pb 4f':
            if MinMaxPb4f == True:                                              #START HERE
                plotMinMax = True                                   
                for i_Max in range (i_start,len(x_BE)+i_start):
                    if x_BE[i_Max] < MaxPb4f: break#                        HERE
                for i_Min in range (i_start,len(x_BE)+i_start):
                    if x_BE[i_Min] < MinPb4f: break#                        HERE       
                plotMax = x_BE[i_Max]
                plotMin = x_BE[i_Min]
                if Div_Area_Check == False: 
                    y_BKG = (y[i_Min-3]+y[i_Min-2]+y[i_Min-1]+y[i_Min])/4       
                    y2 = y - y_BKG                 
                    y2 = 100* y2 / max(y2[i_Max:i_Min-1])             
                    y3 = 100* y3 / max(y3[i_Max:i_Min-1])                       #FINISH HERE    
                else:
                    y_BKG = (y[i_Min-3]+y[i_Min-2]+y[i_Min-1]+y[i_Min])/4       
                    y2 = (y - y_BKG) / Div_Area_Val[j]       
            #df_Pb4f['BE'+Sample_name[j]] = x_BE
            additional = pd.DataFrame({'BE'+Sample_name[j]: x_BE, file[j]:y2})
            new = pd.concat([df_Pb4f,additional],axis =1)
            df_Pb4f = new
            df_Pb4fT.append(Sample_name[j])
            #print(df_Pb4f.head())
            #df_Pb4f[file[j]] = y2
            if (AutoShirley == True) or (len(df1.columns) > j_start + 4):                
                if (AutoShirley == True) or (df1.ix[i_start-2,j_start+4] == 'Backgnd.'):
                    #df_Pb4fs['BE'+Sample_name[j]] = x_BE
                    #df_Pb4fs[file[j]] = y3
                    additional = pd.DataFrame({'BE'+Sample_name[j]: x_BE, file[j]:y3})
                    new = pd.concat([df_Pb4fs,additional],axis =1)
                    df_Pb4fs = new
                    df_Pb4fsT.append(Sample_name[j])
        #----------------------------------------------------

        #-----CHECK CE4D-------------------------------------
        if plotname == 'Ce 4d':
            if MinMaxCe4d == True:                                              #START HERE
                plotMinMax = True                                   
                for i_Max in range (i_start,len(x_BE)+i_start):
                    if x_BE[i_Max] < MaxCe4d: break#                        HERE
                for i_Min in range (i_start,len(x_BE)+i_start):
                    if x_BE[i_Min] < MinCe4d: break#                        HERE       
                plotMax = x_BE[i_Max]
                plotMin = x_BE[i_Min]
                if Div_Area_Check == False: 
                    y_BKG = (y[i_Min-3]+y[i_Min-2]+y[i_Min-1]+y[i_Min])/4       
                    y2 = y - y_BKG                 
                    y2 = 100* y2 / max(y2[i_Max:i_Min-1])             
                    y3 = 100* y3 / max(y3[i_Max:i_Min-1])                       #FINISH HERE    
                else:
                    y_BKG = (y[i_Min-3]+y[i_Min-2]+y[i_Min-1]+y[i_Min])/4       
                    y2 = (y - y_BKG) / Div_Area_Val[j]       
            #df_Ce4d['BE'+Sample_name[j]] = x_BE
            additional = pd.DataFrame({'BE'+Sample_name[j]: x_BE, file[j]:y2})
            new = pd.concat([df_Ce4d,additional],axis =1)
            df_Ce4d = new
            df_Ce4dT.append(Sample_name[j])
            #print(df_Ce4d.head())
            #df_Ce4d[file[j]] = y2
            if (AutoShirley == True) or (len(df1.columns) > j_start + 4):                
                if (AutoShirley == True) or (df1.ix[i_start-2,j_start+4] == 'Backgnd.'):
                    #df_Ce4ds['BE'+Sample_name[j]] = x_BE
                    #df_Ce4ds[file[j]] = y3
                    additional = pd.DataFrame({'BE'+Sample_name[j]: x_BE, file[j]:y3})
                    new = pd.concat([df_Ce4ds,additional],axis =1)
                    df_Ce4ds = new
                    df_Ce4dsT.append(Sample_name[j])
        #----------------------------------------------------

        #-----CHECK CE3D-------------------------------------
        if plotname == 'Ce 3d':
            if MinMaxCe3d == True:                                              #START HERE
                plotMinMax = True                                  
                for i_Max in range (i_start,len(x_BE)+i_start):
                    if x_BE[i_Max] < MaxCe3d: break#                        HERE
                for i_Min in range (i_start,len(x_BE)+i_start):
                    if x_BE[i_Min] < MinCe3d: break#                        HERE       
                plotMax = x_BE[i_Max]
                plotMin = x_BE[i_Min]
                if Div_Area_Check == False:  
                    y_BKG = (y[i_Min-3]+y[i_Min-2]+y[i_Min-1]+y[i_Min])/4       
                    y2 = y - y_BKG                 
                    y2 = 100* y2 / max(y2[i_Max:i_Min-1])             
                    y3 = 100* y3 / max(y3[i_Max:i_Min-1])                       #FINISH HERE            
                else:
                    y_BKG = (y[i_Min-3]+y[i_Min-2]+y[i_Min-1]+y[i_Min])/4       
                    y2 = (y - y_BKG) / Div_Area_Val[j]       
            #df_Ce3d['BE'+Sample_name[j]] = x_BE
            #df_Ce3d[file[j]] = y2
            additional = pd.DataFrame({'BE'+Sample_name[j]: x_BE, file[j]:y2})
            new = pd.concat([df_Ce3d,additional],axis =1)
            df_Ce3d = new            
            df_Ce3dT.append(Sample_name[j])
            if (AutoShirley == True) or (len(df1.columns) > j_start + 4):                
                if (AutoShirley == True) or (df1.ix[i_start-2,j_start+4] == 'Backgnd.'):
                    #df_Ce3ds['BE'+Sample_name[j]] = x_BE
                    #df_Ce3ds[file[j]] = y3
                    additional = pd.DataFrame({'BE'+Sample_name[j]: x_BE, file[j]:y3})
                    new = pd.concat([df_Ce3ds,additional],axis =1)
                    df_Ce3ds = new
                    df_Ce3dsT.append(Sample_name[j])
        #----------------------------------------------------

        #-----CHECK lA5D-------------------------------------
        elif plotname == 'La 5d':
            if MinMaxLa5d == True:                                              #START HERE
                plotMinMax = True                                 
                for i_Max in range (i_start,len(x_BE)+i_start):
                    if x_BE[i_Max] < MaxLa5d: break#                        HERE
                for i_Min in range (i_start,len(x_BE)+i_start):
                    if x_BE[i_Min] < MinLa5d: break#                        HERE       
                if Div_Area_Check == False:    
                    plotMax = x_BE[i_Max]
                    plotMin = x_BE[i_Min]
                    print(i_Max, '   ', i_Min)  #max(y2[i_Max:i_Min-1]))
                    y_BKG = (y[i_Min-3]+y[i_Min-2]+y[i_Min-1]+y[i_Min])/4  
                    y2 = y - y_BKG                 
                    y2 = 100* y2 / max(y2[i_Max:i_Min-1])             
                    y3 = 100* y3 / max(y3[i_Max:i_Min-1])                       #FINISH HERE            
                else:
                    y_BKG = (y[i_Min-3]+y[i_Min-2]+y[i_Min-1]+y[i_Min])/4       
                    y2 = (y - y_BKG) / Div_Area_Val[j]       
            #df_La5d['BE'+Sample_name[j]] = x_BE
            #df_La5d[file[j]] = y2
            additional = pd.DataFrame({'BE'+Sample_name[j]: x_BE, file[j]:y2})
            new = pd.concat([df_La5d,additional],axis =1)
            df_La5d = new
            df_La5dT.append(Sample_name[j])
            if (AutoShirley == True) or (len(df1.columns) > j_start + 4):                
                if (AutoShirley == True) or (df1.ix[i_start-2,j_start+4] == 'Backgnd.'):
                    #df_La5ds['BE'+Sample_name[j]] = x_BE
                    #df_La5ds[file[j]] = y3
                    additional = pd.DataFrame({'BE'+Sample_name[j]: x_BE, file[j]:y3})
                    new = pd.concat([df_La5ds,additional],axis =1)
                    df_La5ds = new
                    df_La5dsT.append(Sample_name[j])
        #----------------------------------------------------

        #-----CHECK lA3D-------------------------------------
        elif plotname == 'La 3d':
            if MinMaxLa3d == True:                                              #START HERE
                plotMinMax = True                                 
                for i_Max in range (i_start,len(x_BE)+i_start):
                    if x_BE[i_Max] < MaxLa3d: break#                        HERE
                for i_Min in range (i_start,len(x_BE)+i_start):
                    if x_BE[i_Min] < MinLa3d: break#                        HERE       
                if Div_Area_Check == False:    
                    plotMax = x_BE[i_Max]
                    plotMin = x_BE[i_Min]
                    y_BKG = (y[i_Min-3]+y[i_Min-2]+y[i_Min-1]+y[i_Min])/4       
                    y2 = y - y_BKG                 
                    y2 = 100* y2 / max(y2[i_Max:i_Min-1])             
                    y3 = 100* y3 / max(y3[i_Max:i_Min-1])                       #FINISH HERE            
                else:
                    y_BKG = (y[i_Min-3]+y[i_Min-2]+y[i_Min-1]+y[i_Min])/4       
                    y2 = (y - y_BKG) / Div_Area_Val[j]       
            #df_La3d['BE'+Sample_name[j]] = x_BE
            #df_La3d[file[j]] = y2
            additional = pd.DataFrame({'BE'+Sample_name[j]: x_BE, file[j]:y2})
            new = pd.concat([df_La3d,additional],axis =1)
            df_La3d = new
            df_La3dT.append(Sample_name[j])
            if (AutoShirley == True) or (len(df1.columns) > j_start + 4):                
                if (AutoShirley == True) or (df1.ix[i_start-2,j_start+4] == 'Backgnd.'):
                    #df_La3ds['BE'+Sample_name[j]] = x_BE
                    #df_La3ds[file[j]] = y3
                    additional = pd.DataFrame({'BE'+Sample_name[j]: x_BE, file[j]:y3})
                    new = pd.concat([df_La3ds,additional],axis =1)
                    df_La3ds = new
                    df_La3dsT.append(Sample_name[j])
        #----------------------------------------------------

        #-----CHECK AG3D-------------------------------------
        elif plotname == 'Ag 3d':
            if MinMaxAg3d == True:                                              #START HERE
                plotMinMax = True                                
                for i_Max in range (i_start,len(x_BE)+i_start):
                    if x_BE[i_Max] < MaxAg3d: break#                        HERE
                for i_Min in range (i_start,len(x_BE)+i_start):
                    if x_BE[i_Min] < MinAg3d: break#                        HERE       
                plotMax = x_BE[i_Max]
                plotMin = x_BE[i_Min]
                if Div_Area_Check == False:    
                    y_BKG = (y[i_Min-3]+y[i_Min-2]+y[i_Min-1]+y[i_Min])/4       
                    y2 = y - y_BKG                 
                    y2 = 100* y2 / max(y2[i_Max:i_Min-1])             
                    y3 = 100* y3 / max(y3[i_Max:i_Min-1])                       #FINISH HERE            
                else:
                    y_BKG = (y[i_Min-3]+y[i_Min-2]+y[i_Min-1]+y[i_Min])/4       
                    y2 = (y - y_BKG) / Div_Area_Val[j]       
            #df_Ag3d['BE'+Sample_name[j]] = x_BE
            #df_Ag3d[file[j]] = y2
            additional = pd.DataFrame({'BE'+Sample_name[j]: x_BE, file[j]:y2})
            new = pd.concat([df_Ag3d,additional],axis =1)
            df_Ag3d = new
            df_Ag3dT.append(Sample_name[j])
            if (AutoShirley == True) or (len(df1.columns) > j_start + 4):                
                if (AutoShirley == True) or (df1.ix[i_start-2,j_start+4] == 'Backgnd.'):
                    #df_Ag3ds['BE'+Sample_name[j]] = x_BE
                    #df_Ag3ds[file[j]] = y3
                    additional = pd.DataFrame({'BE'+Sample_name[j]: x_BE, file[j]:y3})
                    new = pd.concat([df_Ag3ds,additional],axis =1)
                    df_Ag3ds = new
                    df_Ag3dsT.append(Sample_name[j])
        #----------------------------------------------------

        #-----CHECK SR3D-------------------------------------
        elif plotname == 'Sr 3d':
            if MinMaxSr3d == True:                                             #START HERE
                plotMinMax = True                                                  
                for i_Max in range (i_start,len(x_BE)+i_start):
                    if x_BE[i_Max] < MaxSr3d: break#                        HERE
                for i_Min in range (i_start,len(x_BE)+i_start):
                    if x_BE[i_Min] < MinSr3d: break#                        HERE       
                plotMax = x_BE[i_Max]
                plotMin = x_BE[i_Min]
                if Div_Area_Check == False:
                    y_BKG = (y[i_Min-3]+y[i_Min-2]+y[i_Min-1]+y[i_Min])/4       
                    y2 = y - y_BKG                 
                    y2 = 100* y2 / max(y2[i_Max:i_Min-1])
                    y3 = 100* y3 / max(y3[i_Max:i_Min-1])                      #FINISH HERE      
                else:
                    y_BKG = (y[i_Min-3]+y[i_Min-2]+y[i_Min-1]+y[i_Min])/4       
                    y2 = (y - y_BKG) / Div_Area_Val[j]                          
            additional = pd.DataFrame({'BE'+Sample_name[j]: x_BE, file[j]:y2})
            new = pd.concat([df_Sr3d,additional],axis =1)
            df_Sr3d = new
            df_Sr3dT.append(Sample_name[j])
            if (AutoShirley == True) or (len(df1.columns) > j_start + 4):                
                if (AutoShirley == True) or (df1.ix[i_start-2,j_start+4] == 'Backgnd.'):
                    additional = pd.DataFrame({'BE'+Sample_name[j]: x_BE, file[j]:y3})
                    new = pd.concat([df_Sr3ds,additional],axis =1)
                    df_Sr3ds = new
                    df_Sr3dsT.append(Sample_name[j])
        #----------------------------------------------------

        #-----CHECK RU3D-------------------------------------
        elif plotname == 'Ru 3d':
            if MinMaxRu3d == True:                                             #START HERE
                plotMinMax = True                                                  
                for i_Max in range (i_start,len(x_BE)+i_start):
                    if x_BE[i_Max] < MaxRu3d: break#                        HERE
                for i_Min in range (i_start,len(x_BE)+i_start):
                    if x_BE[i_Min] < MinRu3d: break#                        HERE       
                plotMax = x_BE[i_Max]
                plotMin = x_BE[i_Min]
                if Div_Area_Check == False:
                    y_BKG = (y[i_Min-3]+y[i_Min-2]+y[i_Min-1]+y[i_Min])/4       
                    y2 = y - y_BKG                 
                    y2 = 100* y2 / max(y2[i_Max:i_Min-1])
                    y3 = 100* y3 / max(y3[i_Max:i_Min-1])                      #FINISH HERE                    
                else:
                    y_BKG = (y[i_Min-3]+y[i_Min-2]+y[i_Min-1]+y[i_Min])/4       
                    y2 = (y - y_BKG) / Div_Area_Val[j]       
            additional = pd.DataFrame({'BE'+Sample_name[j]: x_BE, file[j]:y2})
            new = pd.concat([df_Ru3d,additional],axis =1)
            df_Ru3d = new
            df_Ru3dT.append(Sample_name[j])
            if (AutoShirley == True) or (len(df1.columns) > j_start + 4):                
                if (AutoShirley == True) or (df1.ix[i_start-2,j_start+4] == 'Backgnd.'):
                    additional = pd.DataFrame({'BE'+Sample_name[j]: x_BE, file[j]:y3})
                    new = pd.concat([df_Ru3ds,additional],axis =1)
                    df_Ru3ds = new
                    df_Ru3dsT.append(Sample_name[j])
        #----------------------------------------------------

        #-----CHECK ZR3D-------------------------------------
        elif plotname == 'Zr 3d':
            if MinMaxZr3d == True:                                              #START HERE
                plotMinMax = True                                   
                for i_Max in range (i_start,len(x_BE)+i_start):
                    if x_BE[i_Max] < MaxZr3d: break#                        HERE
                for i_Min in range (i_start,len(x_BE)+i_start):
                    if x_BE[i_Min] < MinZr3d: break#                        HERE       
                plotMax = x_BE[i_Max]
                plotMin = x_BE[i_Min]
                if Div_Area_Check == False: 
                    y_BKG = (y[i_Min-3]+y[i_Min-2]+y[i_Min-1]+y[i_Min])/4       
                    y2 = y - y_BKG                 
                    y2 = 100* y2 / max(y2[i_Max:i_Min-1])             
                    y3 = 100* y3 / max(y3[i_Max:i_Min-1])                       #FINISH HERE            
                else:
                    y_BKG = (y[i_Min-3]+y[i_Min-2]+y[i_Min-1]+y[i_Min])/4       
                    y2 = (y - y_BKG) / Div_Area_Val[j]       
            #df_Zr3d['BE'+Sample_name[j]] = x_BE
            #df_Zr3d[file[j]] = y2
            additional = pd.DataFrame({'BE'+Sample_name[j]: x_BE, file[j]:y2})
            new = pd.concat([df_Zr3d,additional],axis =1)
            df_Zr3d = new
            df_Zr3dT.append(Sample_name[j])
            if (AutoShirley == True) or (len(df1.columns) > j_start + 4):                
                if (AutoShirley == True) or (df1.ix[i_start-2,j_start+4] == 'Backgnd.'):
                    #df_Zr3ds['BE'+Sample_name[j]] = x_BE
                    #df_Zr3ds[file[j]] = y3
                    additional = pd.DataFrame({'BE'+Sample_name[j]: x_BE, file[j]:y3})
                    new = pd.concat([df_Zr3ds,additional],axis =1)
                    df_Zr3ds = new
                    df_Zr3dsT.append(Sample_name[j])
        #----------------------------------------------------

        #-----CHECK RU3P-------------------------------------
        elif plotname == 'Ru 3p':
            if MinMaxRu3p == True:                                             #START HERE
                plotMinMax = True                                                  
                for i_Max in range (i_start,len(x_BE)+i_start):
                    if x_BE[i_Max] < MaxRu3p: break#                        HERE
                for i_Min in range (i_start,len(x_BE)+i_start):
                    if x_BE[i_Min] < MinRu3p: break#                        HERE       
                plotMax = x_BE[i_Max]
                plotMin = x_BE[i_Min]
                if Div_Area_Check == False:
                    y_BKG = (y[i_Min-3]+y[i_Min-2]+y[i_Min-1]+y[i_Min])/4       
                    y2 = y - y_BKG                 
                    y2 = 100* y2 / max(y2[i_Max:i_Min-1])
                    y3 = 100* y3 / max(y3[i_Max:i_Min-1])                      #FINISH HERE                    
                else:
                    y_BKG = (y[i_Min-3]+y[i_Min-2]+y[i_Min-1]+y[i_Min])/4       
                    y2 = (y - y_BKG) / Div_Area_Val[j]       
            additional = pd.DataFrame({'BE'+Sample_name[j]: x_BE, file[j]:y2})
            new = pd.concat([df_Ru3p,additional],axis =1)
            df_Ru3p = new
            df_Ru3pT.append(Sample_name[j])
            if (AutoShirley == True) or (len(df1.columns) > j_start + 4):                
                if (AutoShirley == True) or (df1.ix[i_start-2,j_start+4] == 'Backgnd.'):
                    additional = pd.DataFrame({'BE'+Sample_name[j]: x_BE, file[j]:y3})
                    new = pd.concat([df_Ru3ps,additional],axis =1)
                    df_Ru3ps = new
                    df_Ru3psT.append(Sample_name[j])
        #----------------------------------------------------

        #-----CHECK SR3P-------------------------------------
        elif plotname == 'Sr 3p':
            if MinMaxSr3p == True:                                              #START HERE
                plotMinMax = True                                  
                for i_Max in range (i_start,len(x_BE)+i_start):
                    if x_BE[i_Max] < MaxSr3p: break#                        HERE
                for i_Min in range (i_start,len(x_BE)+i_start):
                    if x_BE[i_Min] < MinSr3p: break#                        HERE       
                plotMax = x_BE[i_Max]
                plotMin = x_BE[i_Min]
                if Div_Area_Check == False:  
                    y_BKG = (y[i_Min-3]+y[i_Min-2]+y[i_Min-1]+y[i_Min])/4       
                    y2 = y - y_BKG                 
                    y2 = 100* y2 / max(y2[i_Max:i_Min-1])             
                    y3 = 100* y3 / max(y3[i_Max:i_Min-1])                       #FINISH HERE 
                else:
                    y_BKG = (y[i_Min-3]+y[i_Min-2]+y[i_Min-1]+y[i_Min])/4       
                    y2 = (y - y_BKG) / Div_Area_Val[j]       
            additional = pd.DataFrame({'BE'+Sample_name[j]: x_BE, file[j]:y2})
            new = pd.concat([df_Sr3p,additional],axis =1)
            df_Sr3p = new
            df_Sr3pT.append(Sample_name[j])
            if (AutoShirley == True) or (len(df1.columns) > j_start + 4):                
                if (AutoShirley == True) or (df1.ix[i_start-2,j_start+4] == 'Backgnd.'):
                    additional = pd.DataFrame({'BE'+Sample_name[j]: x_BE, file[j]:y3})
                    new = pd.concat([df_Sr3ps,additional],axis =1)
                    df_Sr3ps = new
                    df_Sr3psT.append(Sample_name[j])
        #---------------------------------------------------- 
       
        #-----CHECK CR2P-------------------------------------
        elif plotname == 'Cr 2p':
            additional = pd.DataFrame({'BE'+Sample_name[j]: x_BE, file[j]:y2})
            new = pd.concat([df_Cr2p,additional],axis =1)
            df_Cr2p = new
            df_Cr2pT.append(Sample_name[j])
            if (AutoShirley == True) or (len(df1.columns) > j_start + 4):                
                if (AutoShirley == True) or (df1.ix[i_start-2,j_start+4] == 'Backgnd.'):
                    additional = pd.DataFrame({'BE'+Sample_name[j]: x_BE, file[j]:y3})
                    new = pd.concat([df_Cr2ps,additional],axis =1)
                    df_Cr2ps = new
                    df_Cr2psT.append(Sample_name[j])
        #----------------------------------------------------

        #-----CHECK FE2P-------------------------------------
        elif plotname == 'Fe 2p':
            if MinMaxFe2p == True:                                              #START HERE
                plotMinMax = True                                   
                for i_Max in range (i_start,len(x_BE)+i_start):
                    if x_BE[i_Max] < MaxFe2p: break#                        HERE
                for i_Min in range (i_start,len(x_BE)+i_start):
                    if x_BE[i_Min] < MinFe2p: break#                        HERE       
                plotMax = x_BE[i_Max]
                plotMin = x_BE[i_Min]
                if Div_Area_Check == False: 
                    y_BKG = (y[i_Min-3]+y[i_Min-2]+y[i_Min-1]+y[i_Min])/4       
                    y2 = y - y_BKG                 
                    y2 = 100* y2 / max(y2[i_Max:i_Min-1])             
                    y3 = 100* y3 / max(y3[i_Max:i_Min-1])                       #FINISH HERE            
                else:
                    y_BKG = (y[i_Min-3]+y[i_Min-2]+y[i_Min-1]+y[i_Min])/4       
                    y2 = (y - y_BKG) / Div_Area_Val[j]       
            #df_Fe2p['BE'+Sample_name[j]] = x_BE
            #df_Fe2p[file[j]] = y2
            additional = pd.DataFrame({'BE'+Sample_name[j]: x_BE, file[j]:y2})
            new = pd.concat([df_Fe2p,additional],axis =1)
            df_Fe2p = new
            df_Fe2pT.append(Sample_name[j])
            if (AutoShirley == True) or (len(df1.columns) > j_start + 4):                
                if (AutoShirley == True) or (df1.ix[i_start-2,j_start+4] == 'Backgnd.'):
                    #df_Fe2ps['BE'+Sample_name[j]] = x_BE
                    #df_Fe2ps[file[j]] = y3
                    additional = pd.DataFrame({'BE'+Sample_name[j]: x_BE, file[j]:y3})
                    new = pd.concat([df_Fe2ps,additional],axis =1)
                    df_Fe2ps = new
                    df_Fe2psT.append(Sample_name[j])
        #----------------------------------------------------

        #-----CHECK MN2P-------------------------------------
        elif plotname == 'Mn 2p':
            additional = pd.DataFrame({'BE'+Sample_name[j]: x_BE, file[j]:y2})
            new = pd.concat([df_Mn2p,additional],axis =1)
            df_Mn2p = new
            df_Mn2pT.append(Sample_name[j])
            if (AutoShirley == True) or (len(df1.columns) > j_start + 4):                
                if (AutoShirley == True) or (df1.ix[i_start-2,j_start+4] == 'Backgnd.'):
                    additional = pd.DataFrame({'BE'+Sample_name[j]: x_BE, file[j]:y3})
                    new = pd.concat([df_Mn2ps,additional],axis =1)
                    df_Mn2ps = new
                    df_Mn2psT.append(Sample_name[j])
        #----------------------------------------------------

        #-----CHECK S2P-------------------------------------
        elif plotname == 'S 2p':
            additional = pd.DataFrame({'BE'+Sample_name[j]: x_BE, file[j]:y2})
            new = pd.concat([df_S2p,additional],axis =1)
            df_S2p = new  
            df_S2pT.append(Sample_name[j])
            if (AutoShirley == True) or (len(df1.columns) > j_start + 4):                
                if (AutoShirley == True) or (df1.ix[i_start-2,j_start+4] == 'Backgnd.'):
                    additional = pd.DataFrame({'BE'+Sample_name[j]: x_BE, file[j]:y3})
                    new = pd.concat([df_S2ps,additional],axis =1)
                    df_S2ps = new   
                    df_S2psT.append(Sample_name[j])
        #----------------------------------------------------

        #-----CHECK Al2P-------------------------------------
        elif plotname == 'Al 2p':
            additional = pd.DataFrame({'BE'+Sample_name[j]: x_BE, file[j]:y2})
            new = pd.concat([df_Al2p,additional],axis =1)
            df_Al2p = new  
            df_Al2pT.append(Sample_name[j])
            if (AutoShirley == True) or (len(df1.columns) > j_start + 4):                
                if (AutoShirley == True) or (df1.ix[i_start-2,j_start+4] == 'Backgnd.'):
                    additional = pd.DataFrame({'BE'+Sample_name[j]: x_BE, file[j]:y3})
                    new = pd.concat([df_Al2ps,additional],axis =1)
                    df_Al2ps = new   
                    df_Al2psT.append(Sample_name[j])
        #----------------------------------------------------

        #-----CHECK Si2P-------------------------------------
        elif plotname == 'Si 2p':
            #print(j,'  ',len(y3), '   ', len(x_BE))
            if MinMaxSi2p == True:                                              #START HERE
                plotMinMax = True                                
                for i_Max in range (i_start,len(x_BE)+i_start):
                    if x_BE[i_Max] < MaxSi2p: break#                        HERE
                for i_Min in range (i_start,len(x_BE)+i_start):
                    if x_BE[i_Min] < MinSi2p: break#                        HERE       
                plotMax = x_BE[i_Max]
                plotMin = x_BE[i_Min]
                if Div_Area_Check == False:    
                    y_BKG = (y[i_Min-3]+y[i_Min-2]+y[i_Min-1]+y[i_Min])/4       
                    y2 = y - y_BKG                 
                    y2 = 100* y2 / max(y2[i_Max:i_Min-1])             
                    y3 = 100* y3 / max(y3[i_Max:i_Min-1])                       #FINISH HERE 
                else:
                    y_BKG = (y[i_Min-3]+y[i_Min-2]+y[i_Min-1]+y[i_Min])/4       
                    y2 = (y - y_BKG) / Div_Area_Val[j]       
            additional = pd.DataFrame({'BE'+Sample_name[j]: x_BE, file[j]:y2})
            new = pd.concat([df_Si2p,additional],axis =1)
            df_Si2p = new  
            df_Si2pT.append(Sample_name[j])
            if (AutoShirley == True) or (len(df1.columns) > j_start + 4):                
                if (AutoShirley == True) or (df1.ix[i_start-2,j_start+4] == 'Backgnd.'):
                    additional = pd.DataFrame({'BE'+Sample_name[j]: x_BE, file[j]:y3})
                    new = pd.concat([df_Si2ps,additional],axis =1)
                    df_Si2ps = new  
                    df_Si2psT.append(Sample_name[j])
        #----------------------------------------------------

        #-----CHECK CO2P--------------------------------------
        elif plotname == 'Co 2p':
            #print(j,'  ',len(y3), '   ', len(x_BE))
            if MinMaxCo2p == True:                                              #START HERE
                plotMinMax = True                                
                for i_Max in range (i_start,len(x_BE)+i_start):
                    if x_BE[i_Max] < MaxCo2p: break#                        HERE
                for i_Min in range (i_start,len(x_BE)+i_start):
                    if x_BE[i_Min] < MinCo2p: break#                        HERE       
                plotMax = x_BE[i_Max]
                plotMin = x_BE[i_Min]
                if Div_Area_Check == False:    
                    y_BKG = (y[i_Min-3]+y[i_Min-2]+y[i_Min-1]+y[i_Min])/4       
                    y2 = y - y_BKG                 
                    y2 = 100* y2 / max(y2[i_Max:i_Min-1])             
                    y3 = 100* y3 / max(y3[i_Max:i_Min-1])                       #FINISH HERE 
                else:
                    y_BKG = (y[i_Min-3]+y[i_Min-2]+y[i_Min-1]+y[i_Min])/4       
                    y2 = (y - y_BKG) / Div_Area_Val[j]       
            additional = pd.DataFrame({'BE'+Sample_name[j]: x_BE, file[j]:y2})
            new = pd.concat([df_Co2p,additional],axis =1)
            df_Co2p = new  
            df_Co2pT.append(Sample_name[j])
            if (AutoShirley == True) or (len(df1.columns) > j_start + 4):                
                if (AutoShirley == True) or (df1.ix[i_start-2,j_start+4] == 'Backgnd.'):
                    additional = pd.DataFrame({'BE'+Sample_name[j]: x_BE, file[j]:y3})
                    new = pd.concat([df_Co2ps,additional],axis =1)
                    df_Co2ps = new  
                    df_Co2psT.append(Sample_name[j])
        #----------------------------------------------------

        #-----CHECK Ni2p--------------------------------------
        elif plotname == 'Ni 2p':
            #print(j,'  ',len(y3), '   ', len(x_BE))
            if MinMaxNi2p == True:                                              #START HERE
                plotMinMax = True                                
                for i_Max in range (i_start,len(x_BE)+i_start):
                    if x_BE[i_Max] < MaxNi2p: break#                        HERE
                for i_Min in range (i_start,len(x_BE)+i_start):
                    if x_BE[i_Min] < MinNi2p: break#                        HERE       
                plotMax = x_BE[i_Max]
                plotMin = x_BE[i_Min]
                if Div_Area_Check == False:    
                    y_BKG = (y[i_Min-3]+y[i_Min-2]+y[i_Min-1]+y[i_Min])/4       
                    y2 = y - y_BKG                 
                    y2 = 100* y2 / max(y2[i_Max:i_Min-1])             
                    y3 = 100* y3 / max(y3[i_Max:i_Min-1])                       #FINISH HERE 
                else:
                    y_BKG = (y[i_Min-3]+y[i_Min-2]+y[i_Min-1]+y[i_Min])/4       
                    y2 = (y - y_BKG) / Div_Area_Val[j]                      
            additional = pd.DataFrame({'BE'+Sample_name[j]: x_BE, file[j]:y2})
            new = pd.concat([df_Ni2p,additional],axis =1)
            df_Ni2p = new  
            df_Ni2pT.append(Sample_name[j])
            if (AutoShirley == True) or (len(df1.columns) > j_start + 4):                
                if (AutoShirley == True) or (df1.ix[i_start-2,j_start+4] == 'Backgnd.'):
                    additional = pd.DataFrame({'BE'+Sample_name[j]: x_BE, file[j]:y3})
                    new = pd.concat([df_Ni2ps,additional],axis =1)
                    df_Ni2ps = new  
                    df_Ni2psT.append(Sample_name[j])
        #----------------------------------------------------

        #-----CHECK Ga5p--------------------------------------
        elif plotname == 'Ga 5p':
            #print(j,'  ',len(y3), '   ', len(x_BE))
            if plotname == 'Ga 5p': plotname = 'Ga 2p'
            if MinMaxGa5p == True:                                              #START HERE
                plotMinMax = True                                
                for i_Max in range (i_start,len(x_BE)+i_start):
                    if x_BE[i_Max] < MaxGa5p: break#                        HERE
                for i_Min in range (i_start,len(x_BE)+i_start):
                    if x_BE[i_Min] < MinGa5p: break#                        HERE       
                plotMax = x_BE[i_Max]
                plotMin = x_BE[i_Min]
                if Div_Area_Check == False:    
                    y_BKG = (y[i_Min-3]+y[i_Min-2]+y[i_Min-1]+y[i_Min])/4       
                    y2 = y - y_BKG                 
                    y2 = 100* y2 / max(y2[i_Max:i_Min-1])             
                    y3 = 100* y3 / max(y3[i_Max:i_Min-1])                       #FINISH HERE 
                else:
                    y_BKG = (y[i_Min-3]+y[i_Min-2]+y[i_Min-1]+y[i_Min])/4       
                    y2 = (y - y_BKG) / Div_Area_Val[j]       
            additional = pd.DataFrame({'BE'+Sample_name[j]: x_BE, file[j]:y2})
            new = pd.concat([df_Ga5p,additional],axis =1)
            df_Ga5p = new  
            df_Ga5pT.append(Sample_name[j])
            if (AutoShirley == True) or (len(df1.columns) > j_start + 4):                
                if (AutoShirley == True) or (df1.ix[i_start-2,j_start+4] == 'Backgnd.'):
                    additional = pd.DataFrame({'BE'+Sample_name[j]: x_BE, file[j]:y3})
                    new = pd.concat([df_Ga5ps,additional],axis =1)
                    df_Ga5ps = new  
                    df_Ga5psT.append(Sample_name[j])
        #----------------------------------------------------

        #-----CHECK Ga2p--------------------------------------
        elif plotname == 'Ga 2p':
            #print(j,'  ',len(y3), '   ', len(x_BE))
            if MinMaxGa2p == True:                                              #START HERE
                plotMinMax = True                                
                for i_Max in range (i_start,len(x_BE)+i_start):
                    if x_BE[i_Max] < MaxGa2p: break#                        HERE
                for i_Min in range (i_start,len(x_BE)+i_start):
                    if x_BE[i_Min] < MinGa2p: break#                        HERE       
                plotMax = x_BE[i_Max]
                plotMin = x_BE[i_Min]
                if Div_Area_Check == False:    
                    y_BKG = (y[i_Min-3]+y[i_Min-2]+y[i_Min-1]+y[i_Min])/4       
                    y2 = y - y_BKG                 
                    y2 = 100* y2 / max(y2[i_Max:i_Min-1])             
                    y3 = 100* y3 / max(y3[i_Max:i_Min-1])                       #FINISH HERE 
                else:
                    y_BKG = (y[i_Min-3]+y[i_Min-2]+y[i_Min-1]+y[i_Min])/4       
                    y2 = (y - y_BKG) / Div_Area_Val[j]       
            additional = pd.DataFrame({'BE'+Sample_name[j]: x_BE, file[j]:y2})
            new = pd.concat([df_Ga2p,additional],axis =1)
            df_Ga2p = new  
            df_Ga2pT.append(Sample_name[j])
            if (AutoShirley == True) or (len(df1.columns) > j_start + 4):                
                if (AutoShirley == True) or (df1.ix[i_start-2,j_start+4] == 'Backgnd.'):
                    additional = pd.DataFrame({'BE'+Sample_name[j]: x_BE, file[j]:y3})
                    new = pd.concat([df_Ga2ps,additional],axis =1)
                    df_Ga2ps = new  
                    df_Ga2psT.append(Sample_name[j])
        #----------------------------------------------------

        #-----CHECK P2p--------------------------------------
        elif plotname == 'P 2p':
            #print(j,'  ',len(y3), '   ', len(x_BE))
            if MinMaxP2p == True:                                              #START HERE
                plotMinMax = True                                
                for i_Max in range (i_start,len(x_BE)+i_start):
                    if x_BE[i_Max] < MaxP2p: break#                        HERE
                for i_Min in range (i_start,len(x_BE)+i_start):
                    if x_BE[i_Min] < MinP2p: break#                        HERE       
                plotMax = x_BE[i_Max]
                plotMin = x_BE[i_Min]
                if Div_Area_Check == False:    
                    y_BKG = (y[i_Min-3]+y[i_Min-2]+y[i_Min-1]+y[i_Min])/4       
                    y2 = y - y_BKG                 
                    y2 = 100* y2 / max(y2[i_Max:i_Min-1])             
                    y3 = 100* y3 / max(y3[i_Max:i_Min-1])                       #FINISH HERE 
                else:
                    y_BKG = (y[i_Min-3]+y[i_Min-2]+y[i_Min-1]+y[i_Min])/4       
                    y2 = (y - y_BKG) / Div_Area_Val[j]       
            additional = pd.DataFrame({'BE'+Sample_name[j]: x_BE, file[j]:y2})
            new = pd.concat([df_P2p,additional],axis =1)
            df_P2p = new  
            df_P2pT.append(Sample_name[j])
            if (AutoShirley == True) or (len(df1.columns) > j_start + 4):                
                if (AutoShirley == True) or (df1.ix[i_start-2,j_start+4] == 'Backgnd.'):
                    additional = pd.DataFrame({'BE'+Sample_name[j]: x_BE, file[j]:y3})
                    new = pd.concat([df_P2ps,additional],axis =1)
                    df_P2ps = new  
                    df_P2psT.append(Sample_name[j])
        #----------------------------------------------------

        #-----CHECK O1s--------------------------------------
        elif plotname == 'O 1s':
            #print(j,'  ',len(y3), '   ', len(x_BE))
            if MinMaxO1s == True:                                              #START HERE
                plotMinMax = True                                  
                for i_Max in range (i_start,len(x_BE)+i_start):
                    if x_BE[i_Max] < MaxO1s: break#                        HERE
                for i_Min in range (i_start,len(x_BE)+i_start):
                    if x_BE[i_Min] < MinO1s: break#                        HERE       
                plotMax = x_BE[i_Max]
                plotMin = x_BE[i_Min]
                if Div_Area_Check == False:  
                    y_BKG = (y[i_Min-3]+y[i_Min-2]+y[i_Min-1]+y[i_Min])/4       
                    y2 = y - y_BKG                 
                    y2 = 100* y2 / max(y2[i_Max:i_Min-1])             
                    y3 = 100* y3 / max(y3[i_Max:i_Min-1])                       #FINISH HERE 
                else:
                    y_BKG = (y[i_Min-3]+y[i_Min-2]+y[i_Min-1]+y[i_Min])/4       
                    y2 = (y - y_BKG) / Div_Area_Val[j]      
            additional = pd.DataFrame({'BE'+Sample_name[j]: x_BE, file[j]:y2})
            new = pd.concat([df_O1s,additional],axis =1)
            df_O1s = new  
            df_O1sT.append(Sample_name[j])
            if (AutoShirley == True) or (len(df1.columns) > j_start + 4):                
                if (AutoShirley == True) or (df1.ix[i_start-2,j_start+4] == 'Backgnd.'):
                    additional = pd.DataFrame({'BE'+Sample_name[j]: x_BE, file[j]:y3})
                    new = pd.concat([df_O1ss,additional],axis =1)
                    df_O1ss = new  
                    df_O1ssT.append(Sample_name[j])
        #----------------------------------------------------

        #-----CHECK C1s--------------------------------------
        elif plotname == 'C 1s':
            if MinMaxC1s == True:                                              #START HERE
                plotMinMax = True                              
                for i_Max in range (i_start,len(x_BE)+i_start):
                    if x_BE[i_Max] < MaxC1s: break#                        HERE
                for i_Min in range (i_start,len(x_BE)+i_start):
                    if x_BE[i_Min] < MinC1s: break#                        HERE       
                plotMax = x_BE[i_Max]
                plotMin = x_BE[i_Min]
                if Div_Area_Check == False:      
                    y_BKG = (y[i_Min-3]+y[i_Min-2]+y[i_Min-1]+y[i_Min])/4       
                    y2 = y - y_BKG                 
                    y2 = 100* y2 / max(y2[i_Max:i_Min-1])             
                    y3 = 100* y3 / max(y3[i_Max:i_Min-1])                       #FINISH HERE 
                else:
                    y_BKG = (y[i_Min-3]+y[i_Min-2]+y[i_Min-1]+y[i_Min])/4       
                    y2 = (y - y_BKG) / Div_Area_Val[j]       
            additional = pd.DataFrame({'BE'+Sample_name[j]: x_BE, file[j]:y3})
            new = pd.concat([df_C1s,additional],axis =1)
            df_C1s = new  
            df_C1sT.append(Sample_name[j])
            if (AutoShirley == True) or (len(df1.columns) > j_start + 4):                
                if (AutoShirley == True) or (df1.ix[i_start-2,j_start+4] == 'Backgnd.'):
                    additional = pd.DataFrame({'BE'+Sample_name[j]: x_BE, file[j]:y3})
                    new = pd.concat([df_C1ss,additional],axis =1)
                    df_C1ss = new  
                    df_C1ssT.append(Sample_name[j])
        #----------------------------------------------------
        
        #-----CHECK F1s--------------------------------------
        elif plotname == 'F 1s':
            #print(j,'  ',len(y3), '   ', len(x_BE))
            if MinMaxF1s == True:                                              #START HERE
                plotMinMax = True                                  
                for i_Max in range (i_start,len(x_BE)+i_start):
                    if x_BE[i_Max] < MaxF1s: break#                        HERE
                for i_Min in range (i_start,len(x_BE)+i_start):
                    if x_BE[i_Min] < MinF1s: break#                        HERE       
                plotMax = x_BE[i_Max]
                plotMin = x_BE[i_Min]
                if Div_Area_Check == False:  
                    y_BKG = (y[i_Min-3]+y[i_Min-2]+y[i_Min-1]+y[i_Min])/4       
                    y2 = y - y_BKG                 
                    y2 = 100* y2 / max(y2[i_Max:i_Min-1])             
                    y3 = 100* y3 / max(y3[i_Max:i_Min-1])                       #FINISH HERE 
                else:
                    y_BKG = (y[i_Min-3]+y[i_Min-2]+y[i_Min-1]+y[i_Min])/4       
                    y2 = (y - y_BKG) / Div_Area_Val[j]       
            additional = pd.DataFrame({'BE'+Sample_name[j]: x_BE, file[j]:y2})
            new = pd.concat([df_F1s,additional],axis =1)
            df_F1s = new  
            df_F1sT.append(Sample_name[j])
            if (AutoShirley == True) or (len(df1.columns) > j_start + 4):                
                if (AutoShirley == True) or (df1.ix[i_start-2,j_start+4] == 'Backgnd.'):
                    additional = pd.DataFrame({'BE'+Sample_name[j]: x_BE, file[j]:y3})
                    new = pd.concat([df_F1ss,additional],axis =1)
                    df_F1ss = new  
                    df_F1ssT.append(Sample_name[j])
        #----------------------------------------------------

        #-----CHECK Na1s--------------------------------------
        elif plotname == 'Na 1s':
            #print(j,'  ',len(y3), '   ', len(x_BE))
            if MinMaxNa1s == True:                                              #START HERE
                plotMinMax = True                                  
                for i_Max in range (i_start,len(x_BE)+i_start):
                    if x_BE[i_Max] < MaxNa1s: break#                        HERE
                for i_Min in range (i_start,len(x_BE)+i_start):
                    if x_BE[i_Min] < MinNa1s: break#                        HERE       
                plotMax = x_BE[i_Max]
                plotMin = x_BE[i_Min]
                if Div_Area_Check == False:  
                    y_BKG = (y[i_Min-3]+y[i_Min-2]+y[i_Min-1]+y[i_Min])/4       
                    y2 = y - y_BKG                 
                    y2 = 100* y2 / max(y2[i_Max:i_Min-1])             
                    y3 = 100* y3 / max(y3[i_Max:i_Min-1])                       #FINISH HERE 
                else:
                    y_BKG = (y[i_Min-3]+y[i_Min-2]+y[i_Min-1]+y[i_Min])/4       
                    y2 = (y - y_BKG) / Div_Area_Val[j]       
            additional = pd.DataFrame({'BE'+Sample_name[j]: x_BE, file[j]:y2})
            new = pd.concat([df_Na1s,additional],axis =1)
            df_Na1s = new  
            df_Na1sT.append(Sample_name[j])
            if (AutoShirley == True) or (len(df1.columns) > j_start + 4):                
                if (AutoShirley == True) or (df1.ix[i_start-2,j_start+4] == 'Backgnd.'):
                    additional = pd.DataFrame({'BE'+Sample_name[j]: x_BE, file[j]:y3})
                    new = pd.concat([df_Na1ss,additional],axis =1)
                    df_Na1ss = new  
                    df_Na1ssT.append(Sample_name[j])
        #----------------------------------------------------

        #-----CHECK N1s--------------------------------------
        elif plotname == 'N 1s':
            #print(j,'  ',len(y3), '   ', len(x_BE))
            if MinMaxN1s == True:                                              #START HERE
                #print('HHHHHHHHHHHHHHHHHHHHHHHHHHHHHHHHHHHHHHHHHHH')
                plotMinMax = True                                  
                for i_Max in range (i_start,len(x_BE)+i_start):
                    if x_BE[i_Max] < MaxN1s: break#                        HERE
                for i_Min in range (i_start,len(x_BE)+i_start):
                    if x_BE[i_Min] < MinN1s: break#                        HERE       
                plotMax = x_BE[i_Max]
                plotMin = x_BE[i_Min]
                if Div_Area_Check == False:  
                    y_BKG = (y[i_Min-3]+y[i_Min-2]+y[i_Min-1]+y[i_Min])/4       
                    y2 = y - y_BKG                 
                    y2 = 100* y2 / max(y2[i_Max:i_Min-1])             
                    y3 = 100* y3 / max(y3[i_Max:i_Min-1])                       #FINISH HERE 
                else:
                    y_BKG = (y[i_Min-3]+y[i_Min-2]+y[i_Min-1]+y[i_Min])/4       
                    y2 = (y - y_BKG) / Div_Area_Val[j]       
            additional = pd.DataFrame({'BE'+Sample_name[j]: x_BE, file[j]:y2})
            new = pd.concat([df_N1s,additional],axis =1)
            df_N1s = new  
            df_N1sT.append(Sample_name[j])
            if (AutoShirley == True) or (len(df1.columns) > j_start + 4):                
                if (AutoShirley == True) or (df1.ix[i_start-2,j_start+4] == 'Backgnd.'):
                    additional = pd.DataFrame({'BE'+Sample_name[j]: x_BE, file[j]:y3})
                    new = pd.concat([df_N1ss,additional],axis =1)
                    df_N1ss = new  
                    df_N1ssT.append(Sample_name[j])
        #----------------------------------------------------
                   
        elif plotname == 'Survey':
            df_survey['BE'] = x_BE
            df_survey['Survey'] = y2/max(y2)*100
            additional = pd.DataFrame({'BE'+Sample_name[j]: x_BE, file[j]:y2})
            new = pd.concat([df_survey2,additional],axis =1)
            df_survey2 = new  
            df_surveyT2.append(Sample_name[j])            
            
            
        #-------------PLOT SHIRLEY-------------------------
        #print(plotname, '   ',len(df1.columns), ' >? ', j_start + 4)
        if plotname == 'Survey': a =1
        elif plotname == 'Valence': a =1
        elif (AutoShirley == True) or (len(df1.columns) > j_start + 4):
            if (AutoShirley == True) or (df1.ix[i_start-2,j_start+4] == 'Backgnd.'):
                #print('IN IT IIIIIIIIIIIIIIIIIIIIIIIIIIII')
                plt.figure(2)
                #it was y4 before
                y3xxx = plot_multi(fig_shirley1, x_BE, y3, i5,plotname, plotMinMax, plotMin,plotMax)
                if plotlineS == True:
                    if plotname == 'La 5d': a=1                    
                    else: 
                        #print('plot shirley ', plotname)
                        plot_line(DEF_lines, x_BE,plotMinMax, plotMin,plotMax, y3xxx, 1.12)                
                plt.figure(1)
                i5 = i5 + 1
            elif len(df1.columns) > j_start + 5:
                for i_shirley in range(0, len(df1.columns)):
                    if df1.ix[i_start-2,i_shirley] == 'Backgnd.':    
                        plt.figure(2)
                        #it was y4 before
                        y3xxx = plot_multi(fig_shirley1, x_BE, y3, i5,plotname, plotMinMax, plotMin,plotMax)
                        if plotlineS == True:
                            if plotname == 'La 5d': a=1
                            else: 
                                #print('plot shirley ', plotname)
                                plot_line(DEF_lines, x_BE,plotMinMax, plotMin,plotMax, y3xxx, 1.12)                
                        plt.figure(1)
                        i5 = i5 + 1                        
        #----------------------------------------------------

        #-------------LOOK FOR FITTING-----------------------
        if IsEtch == False:
            IsEtch2 = False
            if len(df1.columns) > j_start + 6:
                IsFit = True                
                plot_fit(x_BE, df1,plotname, 8, plotMinMax, plotMin,plotMax, i_fit)
                plt.figure(1)
                i_fit = i_fit+1
        else:
            IsEtch2 = True
            if len(df1.columns) > j_start + 4:
                if plotname == 'Survey':
                    plot_etch(x_BE, df1,plotname, 10, plotMinMax, plotMin,plotMax, i_etch)
                else :
                    plot_etch(x_BE, df1,plotname, 9, plotMinMax, plotMin,plotMax, i_etch)
                    i_etch = i_etch+1 
                plt.figure(1)
                           
        #----------------------------------------------------

        if plotname == 'Survey':
            

            i1 = i1-1
            #plot_multi(fig, x_BE, y2, i1,plotname, plotMinMax,plotMin,plotMax)

        elif plotname == 'Valence':
            i1 = i1-1
            plt.figure(6)
            #------PLOT CONFIGURATION---------
            plt.axis([x_BE.max()-0, x_BE.min(), 
                      y2.min()-y2.max()*0.05, y2.max()*1.2])
            plt.ylabel("")
            plt.xlabel("Binding Energy (eV)")
            
             
            
    
            #-----REMOVE Y AXIS FROM PLOT-----
            ax = plt.gca()
            ax.yaxis.set_visible(False)  
            #---------------------------------
        
            #-----SET TITLE-------------------
            plt.text(0.01,.92,'Valence band', horizontalalignment = "left", transform=ax.transAxes,
                fontsize=16, fontweight='bold'  )
            #---------------------------------
            plt.plot(x_BE,y2)            
            plt.figure(1)
        else:

            #y5 = 100*y2 / max(y2)  
            y2xx = plot_multi(fig, x_BE, y2, i1,plotname, plotMinMax,plotMin,plotMax)
            if plotline == True:
                if plotname == 'La 5d': a=1
                else:
                    #print('plot ', plotname)
                    plot_line(DEF_lines, x_BE,plotMinMax, plotMin,plotMax, y2xx, 1.12)
        
        i1 = i1+1
        
    
    #-----Deal with Survey plot--------------------------    
    if len(df_survey.columns) == 2:
        plt.figure(3, figsize=(10,5))
        #------PLOT CONFIGURATION---------
        plt.axis([df_survey['BE'].max()-0, df_survey['BE'].min(), 
                  df_survey['Survey'].min()-df_survey['Survey'].max()*0.05, df_survey['Survey'].max()*1.4])
        plt.ylabel("")
        plt.xlabel("Binding Energy (eV)")

        #-----REMOVE Y AXIS FROM PLOT-----
        ax = plt.gca()
        ax.yaxis.set_visible(False)  
        #---------------------------------
    
        #-----SET TITLE-------------------
        plt.text(0.01,.92,'Survey', horizontalalignment = "left", transform=ax.transAxes,
            fontsize=16, fontweight='bold'  )
        #---------------------------------
        
        #-----LABELS PEAKS----------------
        df1 = xl.parse(xl.sheet_names[i3])
        survey_label= df1.ix[1:,1]
        survey_position= df1.ix[1:,3]
        survey_BE = df_survey.ix[i_start:,0]
        survey_y  = df_survey.ix[i_start:,1]
       
        i4 =10
        
        for k in range(1, len(survey_label)+1):
            
            for l in range(i_start, len(survey_BE)):
                i4 = l
                if float(survey_BE[l]) < float(survey_position[k]):
                    i4 = l
                    break
                
            survey_yMAX = survey_y[i4-30:i4+10]
            
            
            ax.annotate(survey_label[k],xy=(survey_position[k], int(10+max(survey_yMAX))),ha="center", va="bottom", 
                             rotation = 90)#, xytext=(survey_position[k], int(10+max(survey_yMAX))))
        
        #---------------------------------
        
        #------PLOT-----------------------
        plt.plot(df_survey['BE'],df_survey['Survey'])#,'k')
        

    #-----END SURVEY PLOT------------------------------------
    
    #-----SAVE PLOT AND DO PRESENTATIONS---------------------
    if IsSaveFig == True:
        
        plt.figure(1)
        plt.savefig(AnaFolder+'PNGs\CoreLevel'+'_'+Sample_name[j]+'.png', dpi = dpiVal, bbox_inches='tight', transparent=False)     
        plt.savefig(AnaFolder+'PNGs\CoreLevel'+'_'+Sample_name[j]+'.svg', bbox_inches='tight', transparent=True)  
        plt.figure(2)
        plt.savefig(AnaFolder+'PNGs\CoreLevelShirley'+'_'+Sample_name[j]+'.png', dpi = dpiVal, bbox_inches='tight', transparent=False)
        plt.savefig(AnaFolder+'PNGs\CoreLevelShirley'+'_'+Sample_name[j]+'.svg', dpi = dpiVal, bbox_inches='tight', transparent=True)
        plt.figure(8)
        plt.savefig(AnaFolder+'PNGs\CoreLevelFit'+'_'+Sample_name[j]+'.png', dpi = dpiVal, bbox_inches='tight', transparent=False)
        plt.savefig(AnaFolder+'PNGs\CoreLevelFit'+'_'+Sample_name[j]+'.svg', dpi = dpiVal, bbox_inches='tight', transparent=False)
        plt.figure(3)
        plt.savefig(AnaFolder+'PNGs\Survey'+'_'+Sample_name[j]+'.png', dpi = dpiVal, bbox_inches='tight', transparent=False)
        plt.savefig(AnaFolder+'PNGs\Survey'+'_'+Sample_name[j]+'.svg', dpi = dpiVal, bbox_inches='tight', transparent=True)
        plt.figure(6)
        plt.savefig(AnaFolder+'PNGs\Valence'+'_'+Sample_name[j]+'.png', dpi = dpiVal, bbox_inches='tight', transparent=False)
        plt.savefig(AnaFolder+'PNGs\Valence'+'_'+Sample_name[j]+'.svg', dpi = dpiVal, bbox_inches='tight', transparent=True)
        plt.figure(9)
        plt.savefig(AnaFolder+'PNGs\CoreLevelEtch'+'_'+Sample_name[j]+'.png', dpi = dpiVal, bbox_inches='tight', transparent=False)
        plt.savefig(AnaFolder+'PNGs\CoreLevelEtch'+'_'+Sample_name[j]+'.svg', dpi = dpiVal, bbox_inches='tight', transparent=True)
        plt.figure(10)
        plt.savefig(AnaFolder+'PNGs\SurveyEtch'+'_'+Sample_name[j]+'.png', dpi = dpiVal, bbox_inches='tight', transparent=False)
        plt.savefig(AnaFolder+'PNGs\SurveyEtch'+'_'+Sample_name[j]+'.svg', dpi = dpiVal, bbox_inches='tight', transparent=True)
        plt.figure(1)
        
        
        #------------CORE LEVELS ON PPTX SLIDE------------------        
        blank_slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(blank_slide_layout)
        shapes = slide.shapes 
        shapes.title.text = Sample_name[j]+ ': Survey'
        pic = slide.shapes.add_picture(AnaFolder+'PNGs\Survey'+'_'+Sample_name[j]+'.png', 0, Cm(3.12), width=Cm(15))
                 
        
        blank_slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(blank_slide_layout)
        shapes = slide.shapes            
        shapes.title.text = Sample_name[j]+ ': Core levels'
        pic = slide.shapes.add_picture(AnaFolder+'PNGs\CoreLevel'+'_'+Sample_name[j]+'.png', 0, Cm(3.12), width=Cm(15))
             
        if IsShirley == True:
            blank_slide_layout = prs.slide_layouts[5]
            slide = prs.slides.add_slide(blank_slide_layout)
            shapes = slide.shapes
            shapes.title.text = Sample_name[j]+ ': Core levels with Shirley correction'
            pic = slide.shapes.add_picture(AnaFolder+'PNGs\CoreLevelShirley'+'_'+Sample_name[j]+'.png', 
                                           0, Cm(3.12), width=Cm(15))

        if IsFit == True:
            blank_slide_layout = prs.slide_layouts[5]
            slide = prs.slides.add_slide(blank_slide_layout)
            shapes = slide.shapes
            shapes.title.text = Sample_name[j]+ ': Peak fitted Core levels'
            pic = slide.shapes.add_picture(AnaFolder+'PNGs\CoreLevelFit'+'_'+Sample_name[j]+'.png', 
                                           0, Cm(3.12), width=Cm(15))

        if IsEtch == True:
            blank_slide_layout = prs.slide_layouts[5]
            slide = prs.slides.add_slide(blank_slide_layout)
            shapes = slide.shapes
            shapes.title.text = Sample_name[j]+ ': Core levels after Etching'
            pic = slide.shapes.add_picture(AnaFolder+'PNGs\CoreLevelEtch'+'_'+Sample_name[j]+'.png', 
                                           0, Cm(3.12), width=Cm(15))   
            
            blank_slide_layout = prs.slide_layouts[5]
            slide = prs.slides.add_slide(blank_slide_layout)
            shapes = slide.shapes
            shapes.title.text = Sample_name[j]+ ': Surveys after Etching'
            pic = slide.shapes.add_picture(AnaFolder+'PNGs\SurveyEtch'+'_'+Sample_name[j]+'.png', 
                                           0, Cm(3.12), width=Cm(15))               

        if IsValence == True:
            blank_slide_layout = prs.slide_layouts[5]
            slide = prs.slides.add_slide(blank_slide_layout)
            shapes = slide.shapes
            shapes.title.text = Sample_name[j]+ ': Valence band'
            pic = slide.shapes.add_picture(AnaFolder+'PNGs\Valence'+'_'+Sample_name[j]+'.png', 
                                           0, Cm(3.12), width=Cm(12))
            
        if isQuant == True:
            print('Quantifying')
            blank_slide_layout = prs.slide_layouts[5]
            slide = prs.slides.add_slide(blank_slide_layout)
            shapes = slide.shapes
            shapes.title.text = Sample_name[j]+ ': Quantifications'
            
            
            Header_Tab = list(Quant_Tab1)
            
            rows = len(Peak_Name)+1
            cols = len(Header_Tab)
            left = Cm(0.5)
            top = Cm(3.5) 
            width = Cm(5.0) + Cm(2.5)*len(Header_Tab)
            height = Cm(0.8)
            
            
            table = shapes.add_table(rows, cols, left, top, width, height).table
            # set column widths
            table.first_col = True
            table.columns[0].width = Cm(5.0)
            for i_tab in range(1, len(Header_Tab)): table.columns[i_tab].width = Cm(2.5)

            
            #print(list(Quant_Tab1))
            for i_tab in range(0, len(Header_Tab)):
                table.cell(0, i_tab).text = Header_Tab[i_tab]
                
                for j_tab in range(0, len(Peak_Name)):
                    '''
                    if i_tab == 1:
                            Quant_Tab1.ix[j_tab,i_tab] = Quant_Tab1.ix[j_tab,i_tab] - CarbonComp[j]+284.8
                            '''
                    table.cell(j_tab+1, i_tab).text = str(Quant_Tab1.ix[j_tab,i_tab])
            
            #table.cell(0, 1).text = Header_Tab[1]
            
                                    
            #df_to_table(slide, Quant_Tab1, 0, Cm(3.12, Cm(12),Cm(10)))
            
            for cell in iter_cells(table):
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(12)
        
            
            rows = len(Header_Tab2)
            cols = 2
            left = Cm(13.5)
            top = Cm(3.5) 
            width = Cm(11.0)
            height = Cm(0.8)
            
            table1 = shapes.add_table(rows, cols, left, top, width, height).table
            table1.first_col = True
            # set column widths
            table1.columns[0].width = Cm(7.0)
            table1.columns[1].width = Cm(4.0)
            
            for i_tab in range(0, len(Header_Tab2)):
                table1.cell(i_tab, 0).text = Quant_Tab3[i_tab][0]      
                table1.cell(i_tab, 1).text = Quant_Tab3[i_tab][1]       

            for cell in iter_cells(table1):
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(12)
            
        prs.save(AnaFolder+P_Fname)
    #------------------------------------------------------            
    #----------------------------------------------------------            
    if UseSourcePython == True: plt.show()
else: plt.close('all')

    
#-----SINGLE PLOT----------------------------------------

# Set the sub figures

fig_multi = plt.figure(4, figsize=(14,14))
fig_multiS = plt.figure(5, figsize=(14,14))
fig_valence = plt.figure(7, figsize=(4,4))
fig_valence2 = plt.figure(11, figsize=(4,4))
fig_survey = plt.figure(10, figsize=(9.3,5))
offset1 = 0.0

if IsHeightAuto == True:
    HeightComp = 1.4
    if len(file) == 2: HeightComp = 1.4
    elif len(file) > 2: HeightComp = 1.6
    elif len(file) > 4: HeightComp = 1.8
    else: HeightComp = 1.4
i2=0
i2s=0

plotlineVal = False

#Val---------------------------------------------------
tmplist = list(df_val)
if len(tmplist)>2:
    plot_comp(df_val, df_valT,7,'VB',True,-2,12.5, 0)
    df_val.to_csv(AnaFolder+'Data_CSV\comp_valence.csv', sep='\t')
    if plotline == True:
        tmplist = list(df_val)
        #print('TMPVAL ' ,len(tmplist))
        print('plot Valence')
        plot_line(DEF_lines, df_val[tmplist[0]], False, 0, 0, df_val[tmplist[1]], 1.12)
    plot_comp(df_val, df_valT,11,'VB',True,Minval, Maxval, 0)
#val---------------------------------------------------


#Ce5d---------------------------------------------------
tmplist = list(df_Ce5d)

#      Checking for labelling----------------
if LabellingAutoHeight == True: 
    TmpDef1 = df_Ce5d
    TmpMin = MinCe5d
    TmpMax = MaxCe5d
    I_maxCOMP = 0
    for i in range(0, int(len(TmpDef1.columns)/2)):
        
        for ix_Max in range (i_start,len(TmpDef1[tmplist[2*i]])+i_start):
            if TmpDef1[tmplist[2*i]].ix[ix_Max] < TmpMax: break
        for ix_Min in range (i_start,len(TmpDef1[tmplist[2*i]])+i_start):
            if TmpDef1[tmplist[2*i]].ix[ix_Min] < TmpMin: break
    
        if I_maxCOMP <  max(TmpDef1[tmplist[2*i+1]].ix[ix_Max:ix_Min]):
            I_maxCOMP = max(TmpDef1[tmplist[2*i+1]].ix[ix_Max:ix_Min])
            LabellingNumber = i
            #print("I_max: ", I_maxCOMP, "FileNumber: ", )
#      Checking for labelling----------------

if len(tmplist)>2 : 
    i2 =plot_comp(df_Ce5d, df_Ce5dT,4,'Ce 3d',MinMaxCe5d, MinCe5d, MaxCe5d,i2)
    df_Ce5d.to_csv(AnaFolder+'Data_CSV\comp_Ce5d.csv', sep='\t')
    if plotline == True:
        #print('TMPCe5d ' ,len(tmplist))
        print('plot Ce 3d')
        plot_line(DEF_lines, df_Ce5d[tmplist[0]], MinMaxCe5d, MinCe5d, MaxCe5d, df_Ce5d[tmplist[2*LabellingNumber+1]], 1.12)

tmplist = list(df_Ce5ds)
if len(tmplist)>2 :
    if IsShirley: 
        i2s =plot_comp(df_Ce5ds, df_Ce5dsT,5,'Ce 3d',MinMaxCe5d, MinCe5d, MaxCe5d,i2s)
        df_Ce5ds.to_csv(AnaFolder+'Data_CSV\comp_Ce5dS.csv', sep='\t')
        if plotline == True: 
            print('plot Ce 3d shirley')
            plot_line(DEF_lines, df_Ce5ds[tmplist[0]], MinMaxCe5d, MinCe5d, MaxCe5d, df_Ce5ds[tmplist[2*LabellingNumber+1]], 1.12)
#Ce5d---------------------------------------------------

#Pb4f---------------------------------------------------
tmplist = list(df_Pb4f)

#      Checking for labelling----------------
if LabellingAutoHeight == True: 
    TmpDef1 = df_Pb4f
    TmpMin = MinPb4f
    TmpMax = MaxPb4f
    I_maxCOMP = 0
    for i in range(0, int(len(TmpDef1.columns)/2)):
        
        for ix_Max in range (i_start,len(TmpDef1[tmplist[2*i]])+i_start):
            if TmpDef1[tmplist[2*i]].ix[ix_Max] < TmpMax: break
        for ix_Min in range (i_start,len(TmpDef1[tmplist[2*i]])+i_start):
            if TmpDef1[tmplist[2*i]].ix[ix_Min] < TmpMin: break
    
        if I_maxCOMP <  max(TmpDef1[tmplist[2*i+1]].ix[ix_Max:ix_Min]):
            I_maxCOMP = max(TmpDef1[tmplist[2*i+1]].ix[ix_Max:ix_Min])
            LabellingNumber = i
            #print("I_max: ", I_maxCOMP, "FileNumber: ", )
#      Checking for labelling----------------

if len(tmplist)>2 : 
    i2 =plot_comp(df_Pb4f, df_Pb4fT,4,'Pb 4f',MinMaxPb4f, MinPb4f, MaxPb4f,i2)
    df_Pb4f.to_csv(AnaFolder+'Data_CSV\comp_Pb4f.csv', sep='\t')
    if plotline == True:
        if UseOffsetComp == True: LabellingNumber = 0
        if len(tmplist)>2 : plot_line(DEF_lines, df_Pb4f[tmplist[0]], MinMaxPb4f, MinPb4f, MaxPb4f, df_Pb4f[tmplist[2*LabellingNumber+1]], 1.12)
tmplist = list(df_Pb4fs)
if len(tmplist)>2 : 
    if IsShirley:
        i2s =plot_comp(df_Pb4fs, df_Pb4fsT,5,'Pb 4f',MinMaxPb4f, MinPb4f, MaxPb4f,i2s)
        df_Pb4fs.to_csv(AnaFolder+'Data_CSV\comp_Pb4fS.csv', sep='\t')
        if plotline == True:
            plot_line(DEF_lines, df_Pb4fs[tmplist[0]], MinMaxPb4f, MinPb4f, MaxPb4f, df_Pb4fs[tmplist[2*LabellingNumber+1]], 1.12)
#Pb4f---------------------------------------------------

#Au4f---------------------------------------------------
tmplist = list(df_Au4f)

#      Checking for labelling----------------
if LabellingAutoHeight == True: 
    TmpDef1 = df_Au4f
    TmpMin = MinAu4f
    TmpMax = MaxAu4f
    I_maxCOMP = 0
    for i in range(0, int(len(TmpDef1.columns)/2)):
        
        for ix_Max in range (i_start,len(TmpDef1[tmplist[2*i]])+i_start):
            if TmpDef1[tmplist[2*i]].ix[ix_Max] < TmpMax: break
        for ix_Min in range (i_start,len(TmpDef1[tmplist[2*i]])+i_start):
            if TmpDef1[tmplist[2*i]].ix[ix_Min] < TmpMin: break
    
        if I_maxCOMP <  max(TmpDef1[tmplist[2*i+1]].ix[ix_Max:ix_Min]):
            I_maxCOMP = max(TmpDef1[tmplist[2*i+1]].ix[ix_Max:ix_Min])
            LabellingNumber = i
            #print("I_max: ", I_maxCOMP, "FileNumber: ", )
#      Checking for labelling----------------

if len(tmplist)>2 : 
    i2 =plot_comp(df_Au4f, df_Au4fT,4,'Au 4f',MinMaxAu4f, MinAu4f, MaxAu4f,i2)
    df_Au4f.to_csv(AnaFolder+'Data_CSV\comp_Au4f.csv', sep='\t')
    if plotline == True:
        if UseOffsetComp == True: LabellingNumber = 0
        if len(tmplist)>2 : plot_line(DEF_lines, df_Au4f[tmplist[0]], MinMaxAu4f, MinAu4f, MaxAu4f, df_Au4f[tmplist[2*LabellingNumber+1]], 1.12)
tmplist = list(df_Au4fs)
if len(tmplist)>2 : 
    if IsShirley:
        i2s =plot_comp(df_Au4fs, df_Au4fsT,5,'Au 4f',MinMaxAu4f, MinAu4f, MaxAu4f,i2s)
        df_Au4fs.to_csv(AnaFolder+'Data_CSV\comp_Au4fS.csv', sep='\t')
        if plotline == True:
            plot_line(DEF_lines, df_Au4fs[tmplist[0]], MinMaxAu4f, MinAu4f, MaxAu4f, df_Au4fs[tmplist[2*LabellingNumber+1]], 1.12)
#Au4f---------------------------------------------------

#Ce4d---------------------------------------------------
tmplist = list(df_Ce4d)

#      Checking for labelling----------------
if LabellingAutoHeight == True: 
    TmpDef1 = df_Ce4d
    TmpMin = MinCe4d
    TmpMax = MaxCe4d
    I_maxCOMP = 0
    for i in range(0, int(len(TmpDef1.columns)/2)):
        
        for ix_Max in range (i_start,len(TmpDef1[tmplist[2*i]])+i_start):
            if TmpDef1[tmplist[2*i]].ix[ix_Max] < TmpMax: break
        for ix_Min in range (i_start,len(TmpDef1[tmplist[2*i]])+i_start):
            if TmpDef1[tmplist[2*i]].ix[ix_Min] < TmpMin: break
    
        if I_maxCOMP <  max(TmpDef1[tmplist[2*i+1]].ix[ix_Max:ix_Min]):
            I_maxCOMP = max(TmpDef1[tmplist[2*i+1]].ix[ix_Max:ix_Min])
            LabellingNumber = i
            #print("I_max: ", I_maxCOMP, "FileNumber: ", )
#      Checking for labelling----------------

if len(tmplist)>2 : 
    i2 =plot_comp(df_Ce4d, df_Ce4dT,4,'Ce 4d',MinMaxCe4d, MinCe4d, MaxCe4d,i2)
    df_Ce4d.to_csv(AnaFolder+'Data_CSV\comp_Ce4d.csv', sep='\t')
    if plotline == True:
        if UseOffsetComp == True: LabellingNumber = 0
        plot_line(DEF_lines, df_Ce4d[tmplist[0]], MinMaxCe4d, MinCe4d, MaxCe4d, df_Ce4d[tmplist[2*LabellingNumber+1]], 1.12)

tmplist = list(df_Ce4ds)
if len(tmplist)>2 : 
    if IsShirley:
        i2s =plot_comp(df_Ce4ds, df_Ce4dsT,5,'Ce 4d',MinMaxCe4d, MinCe4d, MaxCe4d,i2s)
        df_Ce4ds.to_csv(AnaFolder+'Data_CSV\comp_Ce4dS.csv', sep='\t')
        if plotline == True:
            plot_line(DEF_lines, df_Ce4ds[tmplist[0]], MinMaxCe4d, MinCe4d, MaxCe4d, df_Ce4ds[tmplist[2*LabellingNumber+1]], 1.12)
#Ce4d---------------------------------------------------

#La4d---------------------------------------------------
tmplist = list(df_La4d)

#      Checking for labelling----------------
if LabellingAutoHeight == True: 
    TmpDef1 = df_La4d
    TmpMin = MinLa4d
    TmpMax = MaxLa4d
    I_maxCOMP = 0
    for i in range(0, int(len(TmpDef1.columns)/2)):
        
        for ix_Max in range (i_start,len(TmpDef1[tmplist[2*i]])+i_start):
            if TmpDef1[tmplist[2*i]].ix[ix_Max] < TmpMax: break
        for ix_Min in range (i_start,len(TmpDef1[tmplist[2*i]])+i_start):
            if TmpDef1[tmplist[2*i]].ix[ix_Min] < TmpMin: break
    
        if I_maxCOMP <  max(TmpDef1[tmplist[2*i+1]].ix[ix_Max:ix_Min]):
            I_maxCOMP = max(TmpDef1[tmplist[2*i+1]].ix[ix_Max:ix_Min])
            LabellingNumber = i
            #print("I_max: ", I_maxCOMP, "FileNumber: ", )
#      Checking for labelling----------------

if len(tmplist)>2 : 
    i2 =plot_comp(df_La4d, df_La4dT,4,'La 4d',False,0,0,i2)
    df_La4d.to_csv(AnaFolder+'Data_CSV\comp_La4d.csv', sep='\t')
    if plotline == True: 
        if UseOffsetComp == True: LabellingNumber = 0
        plot_line(DEF_lines, df_La4d[tmplist[0]], MinMaxLa4d, MinLa4d, MaxLa4d, df_La4d[tmplist[2*LabellingNumber+1]], 1.12)
tmplist = list(df_La4ds)
if len(tmplist)>2 : 
    if IsShirley:
        i2s =plot_comp(df_La4ds, df_La4dsT,5,'La 4d',False,0,0,i2s)
        df_La4ds.to_csv(AnaFolder+'Data_CSV\comp_La4dS.csv', sep='\t')
        if plotline == True:
            if UseOffsetComp == True: LabellingNumber = 0
            plot_line(DEF_lines, df_La4ds[tmplist[0]], MinMaxLa4d, MinLa4d, MaxLa4d, df_La4ds[tmplist[2*LabellingNumber+1]], 1.12)
#La4d---------------------------------------------------

#Ce3d---------------------------------------------------
tmplist = list(df_Ce3d)

#      Checking for labelling----------------
if LabellingAutoHeight == True: 
    TmpDef1 = df_Ce3d
    TmpMin = MinCe3d
    TmpMax = MaxCe3d
    I_maxCOMP = 0
    for i in range(0, int(len(TmpDef1.columns)/2)):
        
        for ix_Max in range (i_start,len(TmpDef1[tmplist[2*i]])+i_start):
            if TmpDef1[tmplist[2*i]].ix[ix_Max] < TmpMax: break
        for ix_Min in range (i_start,len(TmpDef1[tmplist[2*i]])+i_start):
            if TmpDef1[tmplist[2*i]].ix[ix_Min] < TmpMin: break
    
        if I_maxCOMP <  max(TmpDef1[tmplist[2*i+1]].ix[ix_Max:ix_Min]):
            I_maxCOMP = max(TmpDef1[tmplist[2*i+1]].ix[ix_Max:ix_Min])
            LabellingNumber = i
            #print("I_max: ", I_maxCOMP, "FileNumber: ", )
#      Checking for labelling----------------

if len(tmplist)>2 : 
    i2 =plot_comp(df_Ce3d, df_Ce3dT,4,'Ce 3d',MinMaxCe3d, MinCe3d, MaxCe3d,i2)
    df_Ce3d.to_csv(AnaFolder+'Data_CSV\comp_Ce3d.csv', sep='\t')
    if plotline == True:
        print('plot Ce 3d')
        if UseOffsetComp == True: LabellingNumber = 0
        plot_line(DEF_lines, df_Ce3d[tmplist[0]], MinMaxCe3d, MinCe3d, MaxCe3d, df_Ce3d[tmplist[2*LabellingNumber+1]], 1.12)
tmplist = list(df_Ce3d)
if len(tmplist)>2 : 
    if IsShirley:
        i2s =plot_comp(df_Ce3ds, df_Ce3dsT,5,'Ce 3d',MinMaxCe3d, MinCe3d, MaxCe3d,i2s)
        df_Ce3ds.to_csv(AnaFolder+'Data_CSV\comp_Ce3dS.csv', sep='\t')
        if plotline == True:
            plot_line(DEF_lines, df_Ce3ds[tmplist[0]], MinMaxCe3d, MinCe3d, MaxCe3d, df_Ce3ds[tmplist[2*LabellingNumber+1]], 1.12)
#Ce3d---------------------------------------------------

#La5d---------------------------------------------------
tmplist = list(df_La5d)

#      Checking for labelling----------------
if LabellingAutoHeight == True: 
    TmpDef1 = df_La5d
    TmpMin = MinLa5d
    TmpMax = MaxLa5d
    I_maxCOMP = 0
    for i in range(0, int(len(TmpDef1.columns)/2)):
        
        for ix_Max in range (i_start,len(TmpDef1[tmplist[2*i]])+i_start):
            if TmpDef1[tmplist[2*i]].ix[ix_Max] < TmpMax: break
        for ix_Min in range (i_start,len(TmpDef1[tmplist[2*i]])+i_start):
            if TmpDef1[tmplist[2*i]].ix[ix_Min] < TmpMin: break
    
        if I_maxCOMP <  max(TmpDef1[tmplist[2*i+1]].ix[ix_Max:ix_Min]):
            I_maxCOMP = max(TmpDef1[tmplist[2*i+1]].ix[ix_Max:ix_Min])
            LabellingNumber = i
            #print("I_max: ", I_maxCOMP, "FileNumber: ", )
#      Checking for labelling----------------

if len(tmplist)>2 : 
    i2 =plot_comp(df_La5d, df_La5dT,4,'Ce & La 3d',MinMaxLa5d, MinLa5d, MaxLa5d,i2)
    df_La5d.to_csv(AnaFolder+'Data_CSV\comp_La5d.csv', sep='\t')
    if plotline == True:
        if UseOffsetComp == True: LabellingNumber = 0
        plot_line(DEF_lines, df_La5d[tmplist[0]], MinMaxLa5d, MinLa5d, MaxLa5d, df_La5d[tmplist[2*LabellingNumber+1]], 1.12)
tmplist = list(df_La5ds)
if len(tmplist)>2 : 
    if IsShirley:
        i2s =plot_comp(df_La5ds, df_La5dsT,5,'Ce & La 3d',MinMaxLa5d, MinLa5d, MaxLa5d,i2s)
        df_La5ds.to_csv(AnaFolder+'Data_CSV\comp_La5dS.csv', sep='\t')
        if plotline == True:
            plot_line(DEF_lines, df_La5ds[tmplist[0]], MinMaxLa5d, MinLa5d, MaxLa5d, df_La5ds[tmplist[2*LabellingNumber+1]], 1.12)
#La5d---------------------------------------------------

#La3d---------------------------------------------------
tmplist = list(df_La3d)

#      Checking for labelling----------------
if LabellingAutoHeight == True: 
    TmpDef1 = df_La3d
    TmpMin = MinLa3d
    TmpMax = MaxLa3d
    I_maxCOMP = 0
    for i in range(0, int(len(TmpDef1.columns)/2)):
        
        for ix_Max in range (i_start,len(TmpDef1[tmplist[2*i]])+i_start):
            if TmpDef1[tmplist[2*i]].ix[ix_Max] < TmpMax: break
        for ix_Min in range (i_start,len(TmpDef1[tmplist[2*i]])+i_start):
            if TmpDef1[tmplist[2*i]].ix[ix_Min] < TmpMin: break
    
        if I_maxCOMP <  max(TmpDef1[tmplist[2*i+1]].ix[ix_Max:ix_Min]):
            I_maxCOMP = max(TmpDef1[tmplist[2*i+1]].ix[ix_Max:ix_Min])
            LabellingNumber = i
            #print("I_max: ", I_maxCOMP, "FileNumber: ", )
#      Checking for labelling----------------

if len(tmplist)>2 : 
    i2 =plot_comp(df_La3d, df_La3dT,4,'La 3d',MinMaxLa3d, MinLa3d, MaxLa3d,i2)
    df_La3d.to_csv(AnaFolder+'Data_CSV\comp_La3d.csv', sep='\t')
    if plotline == True:
        if UseOffsetComp == True: LabellingNumber = 0
        plot_line(DEF_lines, df_La3d[tmplist[0]], MinMaxLa3d, MinLa3d, MaxLa3d, df_La3d[tmplist[2*LabellingNumber+1]], 1.12)
tmplist = list(df_La3ds)
if len(tmplist)>2 :
    if IsShirley:
        i2s =plot_comp(df_La3ds, df_La3dsT,5,'La 3d',MinMaxLa3d, MinLa3d, MaxLa3d,i2s)
        df_La3ds.to_csv(AnaFolder+'Data_CSV\comp_La3dS.csv', sep='\t')
        if plotline == True:
            plot_line(DEF_lines, df_La3ds[tmplist[0]], MinMaxLa3d, MinLa3d, MaxLa3d, df_La3ds[tmplist[2*LabellingNumber+1]], 1.12)
#La3d---------------------------------------------------

#Ag3d---------------------------------------------------
i2 =plot_comp(df_Ag3d, df_Ag3dT,4,'Ag 3d',False,0,0,i2)
if plotline == True:
    tmplist = list(df_Ag3d)
    if len(tmplist)>2 : 
        if UseOffsetComp == True: LabellingNumber = 0
        plot_line(DEF_lines, df_Ag3d[tmplist[0]], MinMaxAg3d, MinAg3d, MaxAg3d, df_Ag3d[tmplist[2*LabellingNumber+1]], 1.12)
if IsShirley:
    i2s =plot_comp(df_Ag3ds, df_Ag3dsT,5,'Ag 3d',False,0,0,i2s)
    if plotline == True:
        tmplist = list(df_Ag3ds)
        if len(tmplist)>2 : plot_line(DEF_lines, df_Ag3ds[tmplist[0]], MinMaxAg3d, MinAg3d, MaxAg3d, df_Ag3ds[tmplist[2*LabellingNumber+1]], 1.12)
#Ag3d---------------------------------------------------

'''
#Sr3d---------------------------------------------------
tmplist = list(df_Sr3d)
if len(tmplist)>2 :
    i2 =plot_comp(df_Sr3d, df_Sr3dT,4,'Sr 3d',MinMaxSr3d, MinSr3d, MaxSr3d,i2)
    df_Sr3d.to_csv(AnaFolder+'Data_CSV\comp_Sr3d.csv', sep='\t')
    if plotline == True:
        plot_line(DEF_lines, df_Sr3d[tmplist[0]], MinMaxSr3d, MinSr3d, MaxSr3d, df_Sr3d[tmplist[2*LabellingNumber+1]], 1.12)
tmplist = list(df_Sr3ds)
if len(tmplist)>2 :
    if IsShirley:i2s =plot_comp(df_Sr3ds, df_Sr3dsT,5,'Sr 3d',MinMaxSr3d, MinSr3d, MaxSr3d,i2s)
    df_Sr3ds.to_csv(AnaFolder+'Data_CSV\comp_Sr3dS.csv', sep='\t')
    if plotline == True:
        plot_line(DEF_lines, df_Sr3ds[tmplist[0]], MinMaxSr3d, MinSr3d, MaxSr3d, df_Sr3ds[tmplist[2*LabellingNumber+1]], 1.12)
#Sr3d---------------------------------------------------
'''

#Ru3d---------------------------------------------------
tmplist = list(df_Ru3d)

#      Checking for labelling----------------
if LabellingAutoHeight == True: 
    TmpDef1 = df_Ru3d
    TmpMin = MinRu3d
    TmpMax = MaxRu3d
    I_maxCOMP = 0
    for i in range(0, int(len(TmpDef1.columns)/2)):
        
        for ix_Max in range (i_start,len(TmpDef1[tmplist[2*i]])+i_start):
            if TmpDef1[tmplist[2*i]].ix[ix_Max] < TmpMax: break
        for ix_Min in range (i_start,len(TmpDef1[tmplist[2*i]])+i_start):
            if TmpDef1[tmplist[2*i]].ix[ix_Min] < TmpMin: break
    
        if I_maxCOMP <  max(TmpDef1[tmplist[2*i+1]].ix[ix_Max:ix_Min]):
            I_maxCOMP = max(TmpDef1[tmplist[2*i+1]].ix[ix_Max:ix_Min])
            LabellingNumber = i
            #print("I_max: ", I_maxCOMP, "FileNumber: ", )
#      Checking for labelling----------------

if len(tmplist)>2 : 
    i2 =plot_comp(df_Ru3d, df_Ru3dT,4,'Ru 3d',MinMaxRu3d, MinRu3d, MaxRu3d,i2)
    df_Ru3d.to_csv(AnaFolder+'Data_CSV\comp_Ru3d.csv', sep='\t')
    if plotline == True:
        if UseOffsetComp == True: LabellingNumber = 0
        plot_line(DEF_lines, df_Ru3d[tmplist[0]], MinMaxRu3d, MinRu3d, MaxRu3d, df_Ru3d[tmplist[2*LabellingNumber+1]], 1.12)
tmplist = list(df_Ru3ds)
if len(tmplist)>2 :
    if IsShirley:
        i2s =plot_comp(df_Ru3ds, df_Ru3dsT,5,'Ru 3d',MinMaxRu3d, MinRu3d, MaxRu3d,i2s)
        df_Ru3ds.to_csv(AnaFolder+'Data_CSV\comp_Ru3dS.csv', sep='\t')
        if plotline == True:
            plot_line(DEF_lines, df_Ru3ds[tmplist[0]], MinMaxRu3d, MinRu3d, MaxRu3d, df_Ru3ds[tmplist[2*LabellingNumber+1]], 1.12)
#Ru3d---------------------------------------------------

#Ru3p---------------------------------------------------
tmplist = list(df_Ru3p)

#      Checking for labelling----------------
if LabellingAutoHeight == True: 
    TmpDef1 = df_Ru3p
    TmpMin = MinRu3p
    TmpMax = MaxRu3p
    I_maxCOMP = 0
    for i in range(0, int(len(TmpDef1.columns)/2)):
        
        for ix_Max in range (i_start,len(TmpDef1[tmplist[2*i]])+i_start):
            if TmpDef1[tmplist[2*i]].ix[ix_Max] < TmpMax: break
        for ix_Min in range (i_start,len(TmpDef1[tmplist[2*i]])+i_start):
            if TmpDef1[tmplist[2*i]].ix[ix_Min] < TmpMin: break
    
        if I_maxCOMP <  max(TmpDef1[tmplist[2*i+1]].ix[ix_Max:ix_Min]):
            I_maxCOMP = max(TmpDef1[tmplist[2*i+1]].ix[ix_Max:ix_Min])
            LabellingNumber = i
            #print("I_max: ", I_maxCOMP, "FileNumber: ", )
#      Checking for labelling----------------

if len(tmplist)>2 : 
    i2 =plot_comp(df_Ru3p, df_Ru3pT,4,'Ru 3p',MinMaxRu3p, MinRu3p, MaxRu3p,i2)
    df_Ru3p.to_csv(AnaFolder+'Data_CSV\comp_Ru3p.csv', sep='\t')
    if plotline == True:
        if UseOffsetComp == True: LabellingNumber = 0
        plot_line(DEF_lines, df_Ru3p[tmplist[0]], MinMaxRu3p, MinRu3p, MaxRu3p, df_Ru3p[tmplist[2*LabellingNumber+1]], 1.12)
tmplist = list(df_Ru3ps)
if len(tmplist)>2 :
    if IsShirley:
        i2s =plot_comp(df_Ru3ps, df_Ru3psT,5,'Ru 3p',MinMaxRu3p, MinRu3p, MaxRu3p,i2s)
        df_Ru3ps.to_csv(AnaFolder+'Data_CSV\comp_Ru3pS.csv', sep='\t')
        if plotline == True:
            plot_line(DEF_lines, df_Ru3ps[tmplist[0]], MinMaxRu3p, MinRu3p, MaxRu3p, df_Ru3ps[tmplist[2*LabellingNumber+1]], 1.12)
#Ru3p---------------------------------------------------

#Sr3p---------------------------------------------------
tmplist = list(df_Sr3p)

#      Checking for labelling----------------
if LabellingAutoHeight == True: 
    TmpDef1 = df_Sr3p
    TmpMin = MinSr3p
    TmpMax = MaxSr3p
    I_maxCOMP = 0
    for i in range(0, int(len(TmpDef1.columns)/2)):
        
        for ix_Max in range (i_start,len(TmpDef1[tmplist[2*i]])+i_start):
            if TmpDef1[tmplist[2*i]].ix[ix_Max] < TmpMax: break
        for ix_Min in range (i_start,len(TmpDef1[tmplist[2*i]])+i_start):
            if TmpDef1[tmplist[2*i]].ix[ix_Min] < TmpMin: break
    
        if I_maxCOMP <  max(TmpDef1[tmplist[2*i+1]].ix[ix_Max:ix_Min]):
            I_maxCOMP = max(TmpDef1[tmplist[2*i+1]].ix[ix_Max:ix_Min])
            LabellingNumber = i
            #print("I_max: ", I_maxCOMP, "FileNumber: ", )
#      Checking for labelling----------------

if len(tmplist)>2 : 
    i2 =plot_comp(df_Sr3p, df_Sr3pT,4,'Sr 3p',MinMaxSr3p, MinSr3p, MaxSr3p,i2)
    df_Sr3p.to_csv(AnaFolder+'Data_CSV\comp_Sr3p.csv', sep='\t')
    if plotline == True:
        if UseOffsetComp == True: LabellingNumber = 0
        plot_line(DEF_lines, df_Sr3p[tmplist[0]], MinMaxSr3p, MinSr3p, MaxSr3p, df_Sr3p[tmplist[2*LabellingNumber+1]], 1.12)
tmplist = list(df_Sr3ps)
if len(tmplist)>2 : 
    if IsShirley:
        i2s =plot_comp(df_Sr3ps, df_Sr3psT,5,'Sr 3p',MinMaxSr3p, MinSr3p, MaxSr3p,i2s)
        df_Sr3ps.to_csv(AnaFolder+'Data_CSV\comp_Sr3pS.csv', sep='\t')
        if plotline == True:
            plot_line(DEF_lines, df_Sr3ps[tmplist[0]], MinMaxSr3p, MinSr3p, MaxSr3p, df_Sr3ps[tmplist[2*LabellingNumber+1]], 1.12)
#Sr3p---------------------------------------------------


#Fe2p---------------------------------------------------
tmplist = list(df_Fe2p)

#      Checking for labelling----------------
if LabellingAutoHeight == True: 
    TmpDef1 = df_Fe2p
    TmpMin = MinFe2p
    TmpMax = MaxFe2p
    I_maxCOMP = 0
    for i in range(0, int(len(TmpDef1.columns)/2)):
        
        for ix_Max in range (i_start,len(TmpDef1[tmplist[2*i]])+i_start):
            if TmpDef1[tmplist[2*i]].ix[ix_Max] < TmpMax: break
        for ix_Min in range (i_start,len(TmpDef1[tmplist[2*i]])+i_start):
            if TmpDef1[tmplist[2*i]].ix[ix_Min] < TmpMin: break
    
        if I_maxCOMP <  max(TmpDef1[tmplist[2*i+1]].ix[ix_Max:ix_Min]):
            I_maxCOMP = max(TmpDef1[tmplist[2*i+1]].ix[ix_Max:ix_Min])
            LabellingNumber = i
            #print("I_max: ", I_maxCOMP, "FileNumber: ", )
#      Checking for labelling----------------

if len(tmplist)>2 : 
    i2 =plot_comp(df_Fe2p, df_Fe2pT,4,'Fe 2p',MinMaxFe2p, MinFe2p, MaxFe2p,i2)
    df_Fe2p.to_csv(AnaFolder+'Data_CSV\comp_Fe2p.csv', sep='\t')
    if plotline == True:
        if UseOffsetComp == True: LabellingNumber = 0
        plot_line(DEF_lines, df_Fe2p[tmplist[0]], MinMaxFe2p, MinFe2p, MaxFe2p, df_Fe2p[tmplist[2*LabellingNumber+1]], 1.12)
tmplist = list(df_Fe2ps)
if len(tmplist)>2 : 
    if IsShirley:
        i2s =plot_comp(df_Fe2ps, df_Fe2psT,5,'Fe 2p',MinMaxFe2p, MinFe2p, MaxFe2p,i2s)
        df_Fe2ps.to_csv(AnaFolder+'Data_CSV\comp_Fe2pS.csv', sep='\t')
        if plotline == True:
            plot_line(DEF_lines, df_Fe2ps[tmplist[0]], MinMaxFe2p, MinFe2p, MaxFe2p, df_Fe2ps[tmplist[2*LabellingNumber+1]], 1.12)
#Fe2p---------------------------------------------------


#Mn2p---------------------------------------------------
tmplist = list(df_Mn2p)

#      Checking for labelling----------------
if LabellingAutoHeight == True: 
    TmpDef1 = df_Mn2p
    TmpMin = MinMn2p
    TmpMax = MaxMn2p
    I_maxCOMP = 0
    for i in range(0, int(len(TmpDef1.columns)/2)):
        
        for ix_Max in range (i_start,len(TmpDef1[tmplist[2*i]])+i_start):
            if TmpDef1[tmplist[2*i]].ix[ix_Max] < TmpMax: break
        for ix_Min in range (i_start,len(TmpDef1[tmplist[2*i]])+i_start):
            if TmpDef1[tmplist[2*i]].ix[ix_Min] < TmpMin: break
    
        if I_maxCOMP <  max(TmpDef1[tmplist[2*i+1]].ix[ix_Max:ix_Min]):
            I_maxCOMP = max(TmpDef1[tmplist[2*i+1]].ix[ix_Max:ix_Min])
            LabellingNumber = i
            #print("I_max: ", I_maxCOMP, "FileNumber: ", )
#      Checking for labelling----------------

if len(tmplist)>2 : 
    i2 =plot_comp(df_Mn2p, df_Mn2pT,4,'Mn 2p',False,0,0,i2)
    df_Mn2p.to_csv(AnaFolder+'Data_CSV\comp_Mn2p.csv', sep='\t')
    if plotline == True:
        if UseOffsetComp == True: LabellingNumber = 0
        plot_line(DEF_lines, df_Mn2p[tmplist[0]], MinMaxMn2p, MinMn2p, MaxMn2p, df_Mn2p[tmplist[2*LabellingNumber+1]], 1.12)
tmplist = list(df_Mn2ps)
if len(tmplist)>2 : 
    if IsShirley:
        i2s =plot_comp(df_Mn2ps, df_Mn2psT,5,'Mn 2p',False,0,0,i2s)
        df_Mn2ps.to_csv(AnaFolder+'Data_CSV\comp_Mn2pS.csv', sep='\t')
        if plotline == True:
            plot_line(DEF_lines, df_Mn2ps[tmplist[0]], MinMaxMn2p, MinMn2p, MaxMn2p, df_Mn2ps[tmplist[2*LabellingNumber+1]], 1.12)
#Mn2p---------------------------------------------------


#Cr2p---------------------------------------------------
tmplist = list(df_Cr2p)

#      Checking for labelling----------------
if LabellingAutoHeight == True: 
    TmpDef1 = df_Cr2p
    TmpMin = MinCr2p
    TmpMax = MaxCr2p
    I_maxCOMP = 0
    for i in range(0, int(len(TmpDef1.columns)/2)):
        
        for ix_Max in range (i_start,len(TmpDef1[tmplist[2*i]])+i_start):
            if TmpDef1[tmplist[2*i]].ix[ix_Max] < TmpMax: break
        for ix_Min in range (i_start,len(TmpDef1[tmplist[2*i]])+i_start):
            if TmpDef1[tmplist[2*i]].ix[ix_Min] < TmpMin: break
    
        if I_maxCOMP <  max(TmpDef1[tmplist[2*i+1]].ix[ix_Max:ix_Min]):
            I_maxCOMP = max(TmpDef1[tmplist[2*i+1]].ix[ix_Max:ix_Min])
            LabellingNumber = i
            #print("I_max: ", I_maxCOMP, "FileNumber: ", )
#      Checking for labelling----------------

if len(tmplist)>2 :
    i2 =plot_comp(df_Cr2p, df_Cr2pT,4,'Cr 2p',False,0,0,i2)
    df_Cr2p.to_csv(AnaFolder+'Data_CSV\comp_Cr2p.csv', sep='\t')
    if plotline == True:
        if UseOffsetComp == True: LabellingNumber = 0
        plot_line(DEF_lines, df_Cr2p[tmplist[0]], MinMaxCr2p, MinCr2p, MaxCr2p, df_Cr2p[tmplist[2*LabellingNumber+1]], 1.12)
tmplist = list(df_Cr2ps)
if len(tmplist)>2 : 
    if IsShirley:
        i2s =plot_comp(df_Cr2ps, df_Cr2psT,5,'Cr 2p',False,0,0,i2s)
        df_Cr2ps.to_csv(AnaFolder+'Data_CSV\comp_Cr2pS.csv', sep='\t')
        if plotline == True:
            plot_line(DEF_lines, df_Cr2ps[tmplist[0]], MinMaxCr2p, MinCr2p, MaxCr2p, df_Cr2ps[tmplist[2*LabellingNumber+1]], 1.12)
#Cr2p---------------------------------------------------


#S2p---------------------------------------------------
tmplist = list(df_S2p)

#      Checking for labelling----------------
if LabellingAutoHeight == True: 
    TmpDef1 = df_S2p
    TmpMin = MinS2p
    TmpMax = MaxS2p
    I_maxCOMP = 0
    for i in range(0, int(len(TmpDef1.columns)/2)):
        
        for ix_Max in range (i_start,len(TmpDef1[tmplist[2*i]])+i_start):
            if TmpDef1[tmplist[2*i]].ix[ix_Max] < TmpMax: break
        for ix_Min in range (i_start,len(TmpDef1[tmplist[2*i]])+i_start):
            if TmpDef1[tmplist[2*i]].ix[ix_Min] < TmpMin: break
    
        if I_maxCOMP <  max(TmpDef1[tmplist[2*i+1]].ix[ix_Max:ix_Min]):
            I_maxCOMP = max(TmpDef1[tmplist[2*i+1]].ix[ix_Max:ix_Min])
            LabellingNumber = i
            #print("I_max: ", I_maxCOMP, "FileNumber: ", )
#      Checking for labelling----------------

if len(tmplist)>2 : 
    i2 =plot_comp(df_S2p, df_S2pT,4,'S 2p',False,0,0,i2)
    df_S2p.to_csv(AnaFolder+'Data_CSV\comp_S2p.csv', sep='\t')
    if plotline == True:
        if UseOffsetComp == True: LabellingNumber = 0
        plot_line(DEF_lines, df_S2p[tmplist[0]], MinMaxS2p, MinS2p, MaxS2p, df_S2p[tmplist[2*LabellingNumber+1]], 1.12)
tmplist = list(df_S2ps)
if len(tmplist)>2 : 
    if IsShirley: 
        i2s =plot_comp(df_S2ps, df_S2psT,5,'S 2p',False,0,0,i2s)
        df_S2ps.to_csv(AnaFolder+'Data_CSV\comp_S2pS.csv', sep='\t')
        if plotline == True:
            plot_line(DEF_lines, df_S2ps[tmplist[0]], MinMaxS2p, MinS2p, MaxS2p, df_S2ps[tmplist[2*LabellingNumber+1]], 1.12)
#S2p---------------------------------------------------

#Al2p---------------------------------------------------
tmplist = list(df_Al2p)

#      Checking for labelling----------------
if LabellingAutoHeight == True: 
    TmpDef1 = df_Al2p
    TmpMin = MinAl2p
    TmpMax = MaxAl2p
    I_maxCOMP = 0
    for i in range(0, int(len(TmpDef1.columns)/2)):
        
        for ix_Max in range (i_start,len(TmpDef1[tmplist[2*i]])+i_start):
            if TmpDef1[tmplist[2*i]].ix[ix_Max] < TmpMax: break
        for ix_Min in range (i_start,len(TmpDef1[tmplist[2*i]])+i_start):
            if TmpDef1[tmplist[2*i]].ix[ix_Min] < TmpMin: break
    
        if I_maxCOMP <  max(TmpDef1[tmplist[2*i+1]].ix[ix_Max:ix_Min]):
            I_maxCOMP = max(TmpDef1[tmplist[2*i+1]].ix[ix_Max:ix_Min])
            LabellingNumber = i
            #print("I_max: ", I_maxCOMP, "FileNumber: ", )
#      Checking for labelling----------------

if len(tmplist)>2 : 
    i2 =plot_comp(df_Al2p, df_Al2pT,4,'Al 2p',False,0,0,i2)
    df_Al2p.to_csv(AnaFolder+'Data_CSV\comp_Al2p.csv', sep='\t')
    if plotline == True:
        if UseOffsetComp == True: LabellingNumber = 0
        plot_line(DEF_lines, df_Al2p[tmplist[0]], MinMaxAl2p, MinAl2p, MaxAl2p, df_Al2p[tmplist[2*LabellingNumber+1]], 1.12)
tmplist = list(df_Al2ps)
if len(tmplist)>2 : 
    if IsShirley: 
        i2s =plot_comp(df_Al2ps, df_Al2psT,5,'Al 2p',False,0,0,i2s)
        df_Al2ps.to_csv(AnaFolder+'Data_CSV\comp_Al2pS.csv', sep='\t')
        if plotline == True:
            plot_line(DEF_lines, df_Al2ps[tmplist[0]], MinMaxAl2p, MinAl2p, MaxAl2p, df_Al2ps[tmplist[2*LabellingNumber+1]], 1.12)
#Al2p---------------------------------------------------

#Si2p---------------------------------------------------
tmplist = list(df_Si2p)

#      Checking for labelling----------------
if LabellingAutoHeight == True: 
    TmpDef1 = df_Si2p
    TmpMin = MinSi2p
    TmpMax = MaxSi2p
    I_maxCOMP = 0
    for i in range(0, int(len(TmpDef1.columns)/2)):
        
        for ix_Max in range (i_start,len(TmpDef1[tmplist[2*i]])+i_start):
            if TmpDef1[tmplist[2*i]].ix[ix_Max] < TmpMax: break
        for ix_Min in range (i_start,len(TmpDef1[tmplist[2*i]])+i_start):
            if TmpDef1[tmplist[2*i]].ix[ix_Min] < TmpMin: break
    
        if I_maxCOMP <  max(TmpDef1[tmplist[2*i+1]].ix[ix_Max:ix_Min]):
            I_maxCOMP = max(TmpDef1[tmplist[2*i+1]].ix[ix_Max:ix_Min])
            LabellingNumber = i
            #print("I_max: ", I_maxCOMP, "FileNumber: ", )
#      Checking for labelling----------------

if len(tmplist)>2 : 
    i2 =plot_comp(df_Si2p, df_Si2pT,4,'Si 2p',False,0,0,i2)
    df_Si2p.to_csv(AnaFolder+'Data_CSV\comp_Si2p.csv', sep='\t')
    if plotline == True:
        if UseOffsetComp == True: LabellingNumber = 0
        plot_line(DEF_lines, df_Si2p[tmplist[0]], MinMaxSi2p, MinSi2p, MaxSi2p, df_Si2p[tmplist[2*LabellingNumber+1]], 1.12)
tmplist = list(df_Si2ps)
if len(tmplist)>2 : 
    if IsShirley:
        i2s =plot_comp(df_Si2ps, df_Si2psT,5,'Si 2p',False,0,0,i2s)
        df_Si2ps.to_csv(AnaFolder+'Data_CSV\comp_Si2pS.csv', sep='\t')
        if plotline == True:
            plot_line(DEF_lines, df_Si2ps[tmplist[0]], MinMaxSi2p, MinSi2p, MaxSi2p, df_Si2ps[tmplist[2*LabellingNumber+1]], 1.12)
#Si2p---------------------------------------------------


#Co2p---------------------------------------------------
tmplist = list(df_Co2p)

#      Checking for labelling----------------
if LabellingAutoHeight == True: 
    TmpDef1 = df_Co2p
    TmpMin = MinCo2p
    TmpMax = MaxCo2p
    I_maxCOMP = 0
    for i in range(0, int(len(TmpDef1.columns)/2)):
        
        for ix_Max in range (i_start,len(TmpDef1[tmplist[2*i]])+i_start):
            if TmpDef1[tmplist[2*i]].ix[ix_Max] < TmpMax: break
        for ix_Min in range (i_start,len(TmpDef1[tmplist[2*i]])+i_start):
            if TmpDef1[tmplist[2*i]].ix[ix_Min] < TmpMin: break
    
        if I_maxCOMP <  max(TmpDef1[tmplist[2*i+1]].ix[ix_Max:ix_Min]):
            I_maxCOMP = max(TmpDef1[tmplist[2*i+1]].ix[ix_Max:ix_Min])
            LabellingNumber = i
            #print("I_max: ", I_maxCOMP, "FileNumber: ", )
#      Checking for labelling----------------

if len(tmplist)>2 : 
    i2 = plot_comp(df_Co2p, df_Co2pT,4,'Co 2p',MinMaxCo2p, MinCo2p, MaxCo2p,i2)
    df_Co2p.to_csv(AnaFolder+'Data_CSV\comp_Co2p.csv', sep='\t')
    if plotline == True:
        if UseOffsetComp == True: LabellingNumber = 0
        plot_line(DEF_lines, df_Co2p[tmplist[0]], MinMaxCo2p, MinCo2p, MaxCo2p, df_Co2p[tmplist[2*LabellingNumber+1]], 1.12)
tmplist = list(df_Co2ps)
if len(tmplist)>2 : 
    if IsShirley: 
        i2s =plot_comp(df_Co2ps, df_Co2psT,5,'Co 2p',MinMaxCo2p, MinCo2p, MaxCo2p,i2s)
        df_Co2ps.to_csv(AnaFolder+'Data_CSV\comp_Co2pS.csv', sep='\t')
        if plotline == True:
            plot_line(DEF_lines, df_Co2ps[tmplist[0]], MinMaxCo2p, MinCo2p, MaxCo2p, df_Co2ps[tmplist[2*LabellingNumber+1]], 1.12)
#Co2p---------------------------------------------------


#Ni2p---------------------------------------------------
tmplist = list(df_Ni2p)

#      Checking for labelling----------------
if LabellingAutoHeight == True: 
    TmpDef1 = df_Ni2p
    TmpMin = MinNi2p
    TmpMax = MaxNi2p
    I_maxCOMP = 0
    LabellingNumber = 1
    for i in range(0, int(len(TmpDef1.columns)/2)):
        
        for ix_Max in range (i_start,len(TmpDef1[tmplist[2*i]])+i_start):
            if TmpDef1[tmplist[2*i]].ix[ix_Max] < TmpMax: break
        for ix_Min in range (i_start,len(TmpDef1[tmplist[2*i]])+i_start):
            if TmpDef1[tmplist[2*i]].ix[ix_Min] < TmpMin: break
    
        if I_maxCOMP <  max(TmpDef1[tmplist[2*i+1]].ix[ix_Max:ix_Min]):
            I_maxCOMP = max(TmpDef1[tmplist[2*i+1]].ix[ix_Max:ix_Min])
            LabellingNumber = i
            #print("Above I_max: ", I_maxCOMP, "FileNumber: ", LabellingNumber)
        elif i == 0:
            #print("I_max: ", I_maxCOMP, "FileNumber:  ", LabellingNumber )
            I_maxCOMP = max(TmpDef1[tmplist[2*i+1]].ix[ix_Max:ix_Min])
            LabellingNumber = i
#      Checking for labelling----------------

if len(tmplist)>2 : 
    i2 =plot_comp(df_Ni2p, df_Ni2pT,4,'Ni 2p',MinMaxNi2p, MinNi2p, MaxNi2p,i2)
    df_Ni2p.to_csv(AnaFolder+'Data_CSV\comp_Ni2p.csv', sep='\t')
    if plotline == True:
        if UseOffsetComp == True: LabellingNumber = 0
        plot_line(DEF_lines, df_Ni2p[tmplist[0]], MinMaxNi2p, MinNi2p, MaxNi2p, df_Ni2p[tmplist[2*LabellingNumber+1]], 1.12)
tmplist = list(df_Ni2ps)
if len(tmplist)>2 : 
    if IsShirley:
        i2s =plot_comp(df_Ni2ps, df_Ni2psT,5,'Ni 2p',MinMaxNi2p, MinNi2p, MaxNi2p,i2s)
        df_Ni2ps.to_csv(AnaFolder+'Data_CSV\comp_Ni2pS.csv', sep='\t')
        if plotline == True:
            plot_line(DEF_lines, df_Ni2ps[tmplist[0]], MinMaxNi2p, MinNi2p, MaxNi2p, df_Ni2ps[tmplist[2*LabellingNumber+1]], 1.12)
            
#Ni2p---------------------------------------------------

#Ga5p---------------------------------------------------
tmplist = list(df_Ga5p)

#      Checking for labelling----------------
if LabellingAutoHeight == True: 
    TmpDef1 = df_Ga5p
    TmpMin = MinGa5p
    TmpMax = MaxGa5p
    I_maxCOMP = 0
    for i in range(0, int(len(TmpDef1.columns)/2)):
        
        for ix_Max in range (i_start,len(TmpDef1[tmplist[2*i]])+i_start):
            if TmpDef1[tmplist[2*i]].ix[ix_Max] < TmpMax: break
        for ix_Min in range (i_start,len(TmpDef1[tmplist[2*i]])+i_start):
            if TmpDef1[tmplist[2*i]].ix[ix_Min] < TmpMin: break
    
        if I_maxCOMP <  max(TmpDef1[tmplist[2*i+1]].ix[ix_Max:ix_Min]):
            I_maxCOMP = max(TmpDef1[tmplist[2*i+1]].ix[ix_Max:ix_Min])
            LabellingNumber = i
            #print("I_max: ", I_maxCOMP, "FileNumber: ", )
#      Checking for labelling----------------

if len(tmplist)>2 : 
    i2 =plot_comp(df_Ga5p, df_Ga5pT,4,'Ga 2p',MinMaxGa5p, MinGa5p, MaxGa5p,i2)
    df_Ga5p.to_csv(AnaFolder+'Data_CSV\comp_Ga5p.csv', sep='\t')
    if plotline == True:
        if UseOffsetComp == True: LabellingNumber = 0
        plot_line(DEF_lines, df_Ga5p[tmplist[0]], MinMaxGa5p, MinGa5p, MaxGa5p, df_Ga5p[tmplist[2*LabellingNumber+1]], 1.12)
tmplist = list(df_Ga5ps)
if len(tmplist)>2 : 
    if IsShirley:
        i2s =plot_comp(df_Ga5ps, df_Ga5psT,5,'Ga 2p',MinMaxGa5p, MinGa5p, MaxGa5p,i2s)
        df_Ga5ps.to_csv(AnaFolder+'Data_CSV\comp_Ga5pS.csv', sep='\t')
        if plotline == True:
            plot_line(DEF_lines, df_Ga5ps[tmplist[0]], MinMaxGa5p, MinGa5p, MaxGa5p, df_Ga5ps[tmplist[2*LabellingNumber+1]], 1.12)
#Ga5p---------------------------------------------------

#Ga2p---------------------------------------------------
tmplist = list(df_Ga2p)

#      Checking for labelling----------------
if LabellingAutoHeight == True: 
    TmpDef1 = df_Ga2p
    TmpMin = MinGa2p
    TmpMax = MaxGa2p
    I_maxCOMP = 0
    for i in range(0, int(len(TmpDef1.columns)/2)):
        
        for ix_Max in range (i_start,len(TmpDef1[tmplist[2*i]])+i_start):
            if TmpDef1[tmplist[2*i]].ix[ix_Max] < TmpMax: break
        for ix_Min in range (i_start,len(TmpDef1[tmplist[2*i]])+i_start):
            if TmpDef1[tmplist[2*i]].ix[ix_Min] < TmpMin: break
    
        if I_maxCOMP <  max(TmpDef1[tmplist[2*i+1]].ix[ix_Max:ix_Min]):
            I_maxCOMP = max(TmpDef1[tmplist[2*i+1]].ix[ix_Max:ix_Min])
            LabellingNumber = i
            #print("I_max: ", I_maxCOMP, "FileNumber: ", )
#      Checking for labelling----------------

if len(tmplist)>2 : 
    i2 =plot_comp(df_Ga2p, df_Ga2pT,4,'Ga 2p',MinMaxGa2p, MinGa2p, MaxGa2p,i2)
    df_Ga2p.to_csv(AnaFolder+'Data_CSV\comp_Ga2p.csv', sep='\t')
    if plotline == True:
        if UseOffsetComp == True: LabellingNumber = 0
        plot_line(DEF_lines, df_Ga2p[tmplist[0]], MinMaxGa2p, MinGa2p, MaxGa2p, df_Ga2p[tmplist[2*LabellingNumber+1]], 1.12)
tmplist = list(df_Ga2ps)
if len(tmplist)>2 : 
    if IsShirley:
        i2s =plot_comp(df_Ga2ps, df_Ga2psT,5,'Ga 2p',MinMaxGa2p, MinGa2p, MaxGa2p,i2s)
        df_Ga2ps.to_csv(AnaFolder+'Data_CSV\comp_Ga2pS.csv', sep='\t')
        if plotline == True:
            plot_line(DEF_lines, df_Ga2ps[tmplist[0]], MinMaxGa2p, MinGa2p, MaxGa2p, df_Ga2ps[tmplist[2*LabellingNumber+1]], 1.12)
#Ga2p---------------------------------------------------

#P2p---------------------------------------------------
tmplist = list(df_P2p)

#      Checking for labelling----------------
if LabellingAutoHeight == True: 
    TmpDef1 = df_P2p
    TmpMin = MinP2p
    TmpMax = MaxP2p
    I_maxCOMP = 0
    for i in range(0, int(len(TmpDef1.columns)/2)):
        
        for ix_Max in range (i_start,len(TmpDef1[tmplist[2*i]])+i_start):
            if TmpDef1[tmplist[2*i]].ix[ix_Max] < TmpMax: break
        for ix_Min in range (i_start,len(TmpDef1[tmplist[2*i]])+i_start):
            if TmpDef1[tmplist[2*i]].ix[ix_Min] < TmpMin: break
    
        if I_maxCOMP <  max(TmpDef1[tmplist[2*i+1]].ix[ix_Max:ix_Min]):
            I_maxCOMP = max(TmpDef1[tmplist[2*i+1]].ix[ix_Max:ix_Min])
            LabellingNumber = i
            #print("I_max: ", I_maxCOMP, "FileNumber: ", )
#      Checking for labelling----------------

if len(tmplist)>2 : 
    i2 =plot_comp(df_P2p, df_P2pT,4,'P 2p',MinMaxP2p, MinP2p, MaxP2p,i2)
    df_P2p.to_csv(AnaFolder+'Data_CSV\comp_P2p.csv', sep='\t')
    if plotline == True:
        if UseOffsetComp == True: LabellingNumber = 0
        plot_line(DEF_lines, df_P2p[tmplist[0]], MinMaxP2p, MinP2p, MaxP2p, df_P2p[tmplist[2*LabellingNumber+1]], 1.12)
tmplist = list(df_P2ps)
if len(tmplist)>2 : 
    if IsShirley:
        i2s =plot_comp(df_P2ps, df_P2psT,5,'P 2p',MinMaxP2p, MinP2p, MaxP2p,i2s)
        df_P2ps.to_csv(AnaFolder+'Data_CSV\comp_P2pS.csv', sep='\t')
        if plotline == True:
            plot_line(DEF_lines, df_P2ps[tmplist[0]], MinMaxP2p, MinP2p, MaxP2p, df_P2ps[tmplist[2*LabellingNumber+1]], 1.12)
#P2p---------------------------------------------------

#O1s---------------------------------------------------
tmplist = list(df_O1s)

#      Checking for labelling----------------
if LabellingAutoHeight == True: 
    TmpDef1 = df_O1s
    TmpMin = MinO1s
    TmpMax = MaxO1s
    I_maxCOMP = 0
    for i in range(0, int(len(TmpDef1.columns)/2)):
        
        for ix_Max in range (i_start,len(TmpDef1[tmplist[2*i]])+i_start):
            if TmpDef1[tmplist[2*i]].ix[ix_Max] < TmpMax: break
        for ix_Min in range (i_start,len(TmpDef1[tmplist[2*i]])+i_start):
            if TmpDef1[tmplist[2*i]].ix[ix_Min] < TmpMin: break
    
        if I_maxCOMP <  max(TmpDef1[tmplist[2*i+1]].ix[ix_Max:ix_Min]):
            I_maxCOMP = max(TmpDef1[tmplist[2*i+1]].ix[ix_Max:ix_Min])
            LabellingNumber = i
            print("I_max: ", I_maxCOMP, "FileNumber: ", LabellingNumber )
            #print('GGGGGGGGGGGGGGGGGGGGGGGGGGGGGGGGGGGGGGGGGGGGGGGG')
#      Checking for labelling----------------


if len(tmplist)>2 : 
    i2 =plot_comp(df_O1s, df_O1sT,4,'O 1s',MinMaxO1s, MinO1s, MaxO1s,i2)
    df_O1s.to_csv(AnaFolder+'Data_CSV\comp_O1s.csv', sep='\t')
    if plotline == True:
        if UseOffsetComp == True: LabellingNumber = 0
        plot_line(DEF_lines, df_O1s[tmplist[0]], MinMaxO1s, MinO1s, MaxO1s, df_O1s[tmplist[2*LabellingNumber+1]], 1.12)
tmplist = list(df_O1ss)
if len(tmplist)>2 : 
    if IsShirley:
        i2s =plot_comp(df_O1s, df_O1ssT,5,'O 1s',MinMaxO1s, MinO1s, MaxO1s,i2s)
        df_O1ss.to_csv(AnaFolder+'Data_CSV\comp_O1sS.csv', sep='\t')
        if plotline == True:
            plot_line(DEF_lines, df_O1ss[tmplist[0]], MinMaxO1s, MinO1s, MaxO1s, df_O1ss[tmplist[2*LabellingNumber+1]], 1.12)
#O1s---------------------------------------------------

#Sr3d---------------------------------------------------
tmplist = list(df_Sr3d)

#      Checking for labelling----------------
if LabellingAutoHeight == True: 
    TmpDef1 = df_Sr3d
    TmpMin = MinSr3d
    TmpMax = MaxSr3d
    I_maxCOMP = 0
    for i in range(0, int(len(TmpDef1.columns)/2)):
        
        for ix_Max in range (i_start,len(TmpDef1[tmplist[2*i]])+i_start):
            if TmpDef1[tmplist[2*i]].ix[ix_Max] < TmpMax: break
        for ix_Min in range (i_start,len(TmpDef1[tmplist[2*i]])+i_start):
            if TmpDef1[tmplist[2*i]].ix[ix_Min] < TmpMin: break
    
        if I_maxCOMP <  max(TmpDef1[tmplist[2*i+1]].ix[ix_Max:ix_Min]):
            I_maxCOMP = max(TmpDef1[tmplist[2*i+1]].ix[ix_Max:ix_Min])
            LabellingNumber = i
            #print("I_max: ", I_maxCOMP, "FileNumber: ", )
#      Checking for labelling----------------

if len(tmplist)>2 :
    i2 =plot_comp(df_Sr3d, df_Sr3dT,4,'Sr 3d',MinMaxSr3d, MinSr3d, MaxSr3d,i2)
    df_Sr3d.to_csv(AnaFolder+'Data_CSV\comp_Sr3d.csv', sep='\t')
    if plotline == True:        
        if UseOffsetComp == True: LabellingNumber = 0
        plot_line(DEF_lines, df_Sr3d[tmplist[0]], MinMaxSr3d, MinSr3d, MaxSr3d, df_Sr3d[tmplist[2*LabellingNumber+1]], 1.12)
tmplist = list(df_Sr3ds)
if len(tmplist)>2 :
    if IsShirley:
        i2s =plot_comp(df_Sr3ds, df_Sr3dsT,5,'Sr 3d',MinMaxSr3d, MinSr3d, MaxSr3d,i2s)
        df_Sr3ds.to_csv(AnaFolder+'Data_CSV\comp_Sr3dS.csv', sep='\t')
        if plotline == True:
            plot_line(DEF_lines, df_Sr3ds[tmplist[0]], MinMaxSr3d, MinSr3d, MaxSr3d, df_Sr3ds[tmplist[2*LabellingNumber+1]], 1.12)
#Sr3d---------------------------------------------------

#C1s---------------------------------------------------
tmplist = list(df_C1s)

#      Checking for labelling----------------
if LabellingAutoHeight == True: 
    TmpDef1 = df_C1s
    TmpMin = MinC1s
    TmpMax = MaxC1s
    I_maxCOMP = 0
    for i in range(0, int(len(TmpDef1.columns)/2)):
        
        for ix_Max in range (i_start,len(TmpDef1[tmplist[2*i]])+i_start):
            if TmpDef1[tmplist[2*i]].ix[ix_Max] < TmpMax: break
        for ix_Min in range (i_start,len(TmpDef1[tmplist[2*i]])+i_start):
            if TmpDef1[tmplist[2*i]].ix[ix_Min] < TmpMin: break
    
        if I_maxCOMP <  max(TmpDef1[tmplist[2*i+1]].ix[ix_Max:ix_Min]):
            I_maxCOMP = max(TmpDef1[tmplist[2*i+1]].ix[ix_Max:ix_Min])
            LabellingNumber = i
            #print("I_max: ", I_maxCOMP, "FileNumber: ", )
#      Checking for labelling----------------

if len(tmplist)>2 : 
    i2 =plot_comp(df_C1s, df_C1sT,4,'C 1s',MinMaxC1s, MinC1s, MaxC1s,i2)
    df_C1s.to_csv(AnaFolder+'Data_CSV\comp_C1s.csv', sep='\t')
    if plotline == True:
        if UseOffsetComp == True: LabellingNumber = 0
        plot_line(DEF_lines, df_C1s[tmplist[0]], MinMaxC1s, MinC1s, MaxC1s, df_C1s[tmplist[2*LabellingNumber+1]], 1.12)
tmplist = list(df_C1ss)
if len(tmplist)>2 : 
    if IsShirley:
        i2s =plot_comp(df_C1s, df_C1ssT,5,'C 1s',MinMaxC1s, MinC1s, MaxC1s,i2s)
        df_C1ss.to_csv(AnaFolder+'Data_CSV\comp_C1sS.csv', sep='\t')
        if plotline == True:
            plot_line(DEF_lines, df_C1ss[tmplist[0]], MinMaxC1s, MinC1s, MaxC1s, df_C1ss[tmplist[2*LabellingNumber+1]], 1.12)
#C1s---------------------------------------------------


#F1s---------------------------------------------------
tmplist = list(df_F1s)

#      Checking for labelling----------------
if LabellingAutoHeight == True: 
    TmpDef1 = df_F1s
    TmpMin = MinF1s
    TmpMax = MaxF1s
    I_maxCOMP = 0
    for i in range(0, int(len(TmpDef1.columns)/2)):
        
        for ix_Max in range (i_start,len(TmpDef1[tmplist[2*i]])+i_start):
            if TmpDef1[tmplist[2*i]].ix[ix_Max] < TmpMax: break
        for ix_Min in range (i_start,len(TmpDef1[tmplist[2*i]])+i_start):
            if TmpDef1[tmplist[2*i]].ix[ix_Min] < TmpMin: break
    
        if I_maxCOMP <  max(TmpDef1[tmplist[2*i+1]].ix[ix_Max:ix_Min]):
            I_maxCOMP = max(TmpDef1[tmplist[2*i+1]].ix[ix_Max:ix_Min])
            LabellingNumber = i
            #print("I_max: ", I_maxCOMP, "FileNumber: ", )
#      Checking for labelling----------------

if len(tmplist)>2 : 
    i2 =plot_comp(df_F1s, df_F1sT,4,'F 1s',False,0,0,i2)
    df_F1s.to_csv(AnaFolder+'Data_CSV\comp_F1s.csv', sep='\t')
    if plotline == True:
        if UseOffsetComp == True: LabellingNumber = 0
        plot_line(DEF_lines, df_F1s[tmplist[0]], MinMaxF1s, MinF1s, MaxF1s, df_F1s[tmplist[2*LabellingNumber+1]], 1.12)
tmplist = list(df_F1ss)
if len(tmplist)>2 : 
    if IsShirley:
        i2s =plot_comp(df_F1s, df_F1ssT,5,'F 1s',False,0,0,i2s)
        df_F1ss.to_csv(AnaFolder+'Data_CSV\comp_F1sS.csv', sep='\t')
        if plotline == True:
            plot_line(DEF_lines, df_F1ss[tmplist[0]], MinMaxF1s, MinF1s, MaxF1s, df_F1ss[tmplist[2*LabellingNumber+1]], 1.12)
#F1s---------------------------------------------------

#Na1s---------------------------------------------------
tmplist = list(df_Na1s)

#      Checking for labelling----------------
if LabellingAutoHeight == True: 
    TmpDef1 = df_Na1s
    TmpMin = MinNa1s
    TmpMax = MaxNa1s
    I_maxCOMP = 0
    for i in range(0, int(len(TmpDef1.columns)/2)):
        
        for ix_Max in range (i_start,len(TmpDef1[tmplist[2*i]])+i_start):
            if TmpDef1[tmplist[2*i]].ix[ix_Max] < TmpMax: break
        for ix_Min in range (i_start,len(TmpDef1[tmplist[2*i]])+i_start):
            if TmpDef1[tmplist[2*i]].ix[ix_Min] < TmpMin: break
    
        if I_maxCOMP <  max(TmpDef1[tmplist[2*i+1]].ix[ix_Max:ix_Min]):
            I_maxCOMP = max(TmpDef1[tmplist[2*i+1]].ix[ix_Max:ix_Min])
            LabellingNumber = i
            #print("I_max: ", I_maxCOMP, "FileNumber: ", )
#      Checking for labelling----------------

if len(tmplist)>2 : 
    i2 =plot_comp(df_Na1s, df_Na1sT,4,'Na 1s',False,0,0,i2)
    df_Na1s.to_csv(AnaFolder+'Data_CSV\comp_Na1s.csv', sep='\t')
    if plotline == True:
        if UseOffsetComp == True: LabellingNumber = 0
        plot_line(DEF_lines, df_Na1s[tmplist[0]], MinMaxNa1s, MinNa1s, MaxNa1s, df_Na1s[tmplist[2*LabellingNumber+1]], 1.12)
tmplist = list(df_Na1ss)
if len(tmplist)>2 : 
    if IsShirley:
        i2s =plot_comp(df_Na1s, df_Na1ssT,5,'Na 1s',False,0,0,i2s)
        df_Na1ss.to_csv(AnaFolder+'Data_CSV\comp_Na1sS.csv', sep='\t')
        if plotline == True:
            plot_line(DEF_lines, df_Na1ss[tmplist[0]], MinMaxNa1s, MinNa1s, MaxNa1s, df_Na1ss[tmplist[2*LabellingNumber+1]], 1.12)
#Na1s---------------------------------------------------

#N1s---------------------------------------------------
tmplist = list(df_N1s)

#      Checking for labelling----------------
if LabellingAutoHeight == True: 
    TmpDef1 = df_N1s
    TmpMin = MinN1s
    TmpMax = MaxN1s
    I_maxCOMP = 0
    for i in range(0, int(len(TmpDef1.columns)/2)):
        
        for ix_Max in range (i_start,len(TmpDef1[tmplist[2*i]])+i_start):
            if TmpDef1[tmplist[2*i]].ix[ix_Max] < TmpMax: break
        for ix_Min in range (i_start,len(TmpDef1[tmplist[2*i]])+i_start):
            if TmpDef1[tmplist[2*i]].ix[ix_Min] < TmpMin: break
    
        if I_maxCOMP <  max(TmpDef1[tmplist[2*i+1]].ix[ix_Max:ix_Min]):
            I_maxCOMP = max(TmpDef1[tmplist[2*i+1]].ix[ix_Max:ix_Min])
            LabellingNumber = i
            #print("I_max: ", I_maxCOMP, "FileNumber: ", )
#      Checking for labelling----------------

if len(tmplist)>2 : 
    i2 =plot_comp(df_N1s, df_N1sT,4,'N 1s',False,0,0,i2)
    df_N1s.to_csv(AnaFolder+'Data_CSV\comp_N1s.csv', sep='\t')
    if plotline == True:
        if UseOffsetComp == True: LabellingNumber = 0
        plot_line(DEF_lines, df_N1s[tmplist[0]], MinMaxN1s, MinN1s, MaxN1s, df_N1s[tmplist[2*LabellingNumber+1]], 1.12)
tmplist = list(df_N1ss)
if len(tmplist)>2 : 
    if IsShirley:
        i2s =plot_comp(df_N1s, df_N1ssT,5,'N 1s',False,0,0,i2s)
        df_N1ss.to_csv(AnaFolder+'Data_CSV\comp_N1sS.csv', sep='\t')
        if plotline == True:
            plot_line(DEF_lines, df_N1ss[tmplist[0]], MinMaxN1s, MinN1s, MaxN1s, df_N1ss[tmplist[2*LabellingNumber+1]], 1.12)
#N1s---------------------------------------------------

#Zr3d---------------------------------------------------
tmplist = list(df_Zr3d)

#      Checking for labelling----------------
if LabellingAutoHeight == True: 
    TmpDef1 = df_Zr3d
    TmpMin = MinZr3d
    TmpMax = MaxZr3d
    I_maxCOMP = 0
    for i in range(0, int(len(TmpDef1.columns)/2)):
        
        for ix_Max in range (i_start,len(TmpDef1[tmplist[2*i]])+i_start):
            if TmpDef1[tmplist[2*i]].ix[ix_Max] < TmpMax: break
        for ix_Min in range (i_start,len(TmpDef1[tmplist[2*i]])+i_start):
            if TmpDef1[tmplist[2*i]].ix[ix_Min] < TmpMin: break
    
        if I_maxCOMP <  max(TmpDef1[tmplist[2*i+1]].ix[ix_Max:ix_Min]):
            I_maxCOMP = max(TmpDef1[tmplist[2*i+1]].ix[ix_Max:ix_Min])
            LabellingNumber = i
            #print("I_max: ", I_maxCOMP, "FileNumber: ", )
#      Checking for labelling----------------

if len(tmplist)>2 : 
    i2 =plot_comp(df_Zr3d, df_Zr3dT,4,'Zr 3d',MinMaxZr3d, MinZr3d, MaxZr3d,i2)
    df_Zr3d.to_csv(AnaFolder+'Data_CSV\comp_Zr3d.csv', sep='\t')
    if plotline == True:
        if UseOffsetComp == True: LabellingNumber = 0
        plot_line(DEF_lines, df_Zr3d[tmplist[0]], MinMaxZr3d, MinZr3d, MaxZr3d, df_Zr3d[tmplist[2*LabellingNumber+1]], 1.12)
tmplist = list(df_Zr3ds)
if len(tmplist)>2 :
    if IsShirley:
        i2s =plot_comp(df_Zr3ds, df_Zr3dsT,5,'Zr 3d',MinMaxZr3d, MinZr3d, MaxZr3d,i2s)
        df_Zr3ds.to_csv(AnaFolder+'Data_CSV\comp_Zr3dS.csv', sep='\t')
        if plotline == True:
            plot_line(DEF_lines, df_Zr3ds[tmplist[0]], MinMaxZr3d, MinZr3d, MaxZr3d, df_Zr3ds[tmplist[2*LabellingNumber+1]], 1.12)
#Zr3d---------------------------------------------------

plt.figure(10)
plot_comp(df_survey2, df_surveyT2,10,'Survey',False,0,0, 0)

if i2 >0: IsComp = True
if i2s >0: IsComp = True

#-----SAVE FIGS AND DO PRESENTATION------------------------------
if IsSaveFig == True:
    #if IsComp == True:
    plt.savefig(AnaFolder+'PNGs\Survey_Comp'+'.png', dpi = dpiVal, bbox_inches='tight', transparent=False) 
    plt.savefig(AnaFolder+'PNGs\Survey_Comp'+'.svg', dpi = dpiVal, bbox_inches='tight', transparent=True) 
    plt.figure(4)
    plt.savefig(AnaFolder+'PNGs\CoreLevel_Comp'+'.png', dpi = dpiVal, bbox_inches='tight', transparent=False)
    plt.savefig(AnaFolder+'PNGs\CoreLevel_Comp'+'.svg', dpi = dpiVal, bbox_inches='tight', transparent=True)
    plt.figure(5)
    plt.savefig(AnaFolder+'PNGs\CoreLevel_Comp_Shirley'+'.png', dpi = dpiVal, bbox_inches='tight', transparent=False)
    plt.savefig(AnaFolder+'PNGs\CoreLevel_Comp_Shirley'+'.svg', dpi = dpiVal, bbox_inches='tight', transparent=True)
    plt.figure(7)
    plt.savefig(AnaFolder+'PNGs\Valence_Comp'+'.png', dpi = dpiVal, bbox_inches='tight', transparent=False)
    plt.savefig(AnaFolder+'PNGs\Valence_Comp'+'.svg', dpi = dpiVal, bbox_inches='tight', transparent=True)
    plt.figure(11)
    plt.savefig(AnaFolder+'PNGs\Valence_Comp2'+'.png', dpi = dpiVal, bbox_inches='tight', transparent=False)
    plt.savefig(AnaFolder+'PNGs\Valence_Comp2'+'.svg', dpi = dpiVal, bbox_inches='tight', transparent=True)
    
    #------------CORE LEVELS COMP ON PPTX SLIDE------------------
    if IsPresentation == True:            
        if IsComp == True:

            blank_slide_layout = prs.slide_layouts[5]
            slide = prs.slides.add_slide(blank_slide_layout)
            shapes = slide.shapes 
            shapes.title.text = 'Comparisons: Survey'
            pic = slide.shapes.add_picture(AnaFolder+'PNGs\Survey_Comp'+'.png', 0, Cm(3.12), width=Cm(15))  
            prs.save(AnaFolder+P_Fname)            
            
            blank_slide_layout = prs.slide_layouts[5]
            slide = prs.slides.add_slide(blank_slide_layout)
            shapes = slide.shapes 
            shapes.title.text = 'Comparisons: Core levels'
            pic = slide.shapes.add_picture(AnaFolder+'PNGs\CoreLevel_Comp'+'.png', 0, Cm(3.12), width=Cm(15)) 
            
            blank_slide_layout = prs.slide_layouts[5]
            slide = prs.slides.add_slide(blank_slide_layout)
            shapes = slide.shapes 
            shapes.title.text = 'Comparisons: Core levels with Shirley correction'
            pic = slide.shapes.add_picture(AnaFolder+'PNGs\CoreLevel_Comp_Shirley'+'.png', 0, Cm(3.12), width=Cm(15))    
            prs.save(AnaFolder+P_Fname)
            
            blank_slide_layout = prs.slide_layouts[5]
            slide = prs.slides.add_slide(blank_slide_layout)
            shapes = slide.shapes 
            shapes.title.text = 'Comparisons: Valence band'
            pic = slide.shapes.add_picture(AnaFolder+'PNGs\Valence_Comp'+'.png', 0, Cm(3.12), height=Cm(12))   
            pic = slide.shapes.add_picture(AnaFolder+'PNGs\Valence_Comp2'+'.png', Cm(13), Cm(3.12), height=Cm(12))
            prs.save(AnaFolder+P_Fname)
    
    #-----------------------------------------------------------
#---------------------------------------------------------------


if UseSourcePython == True: 
    plt.show()

for i in range(1,10) : plt.figure(i).clf()
plt.close('all')


